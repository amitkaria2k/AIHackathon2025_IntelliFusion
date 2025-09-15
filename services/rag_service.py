"""RAG (Retrieval-Augmented Generation) service for project data"""

import os
import hashlib
import tempfile
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import sqlite3
import json

# Vector and embedding imports
try:
    from sentence_transformers import SentenceTransformer
    import numpy as np
    from sklearn.metrics.pairwise import cosine_similarity
    HAS_RAG_DEPENDENCIES = True
except ImportError:
    HAS_RAG_DEPENDENCIES = False

# File processing imports
import PyPDF2
import docx
from pptx import Presentation
import pandas as pd

class RAGService:
    """Service for handling RAG operations with project data"""
    
    def __init__(self, db_manager, model_name: str = "all-MiniLM-L6-v2"):
        self.db_manager = db_manager
        self.model_name = model_name
        self.model = None
        self.chunk_size = 1000
        self.chunk_overlap = 200
        
        if HAS_RAG_DEPENDENCIES:
            try:
                self.model = SentenceTransformer(model_name)
            except Exception as e:
                print(f"Warning: Could not load embedding model: {e}")
                self.model = None
    
    def is_available(self) -> bool:
        """Check if RAG service is available"""
        return HAS_RAG_DEPENDENCIES and self.model is not None
    
    def extract_text_from_file(self, file_obj, filename: str) -> str:
        """Extract text content from various file types"""
        try:
            file_extension = filename.lower().split('.')[-1]
            
            if file_extension == 'pdf':
                pdf_reader = PyPDF2.PdfReader(file_obj)
                text = ""
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
                return text
                
            elif file_extension == 'docx':
                doc = docx.Document(file_obj)
                text = ""
                for paragraph in doc.paragraphs:
                    text += paragraph.text + "\n"
                # Extract tables
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join([cell.text for cell in row.cells])
                        text += row_text + "\n"
                return text
                
            elif file_extension in ['xlsx', 'xls']:
                excel_data = pd.read_excel(file_obj, sheet_name=None)
                text = ""
                for sheet_name, df in excel_data.items():
                    text += f"\n--- Sheet: {sheet_name} ---\n"
                    text += df.to_string(index=False) + "\n"
                return text
                
            elif file_extension == 'pptx':
                prs = Presentation(file_obj)
                text = ""
                for slide_num, slide in enumerate(prs.slides, 1):
                    text += f"\n--- Slide {slide_num} ---\n"
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
                
            elif file_extension in ['txt', 'md']:
                return file_obj.read().decode('utf-8')
                
            elif file_extension == 'json':
                data = json.loads(file_obj.read().decode('utf-8'))
                return json.dumps(data, indent=2)
                
            else:
                # For unknown file types, try to read as text
                try:
                    return file_obj.read().decode('utf-8', errors='ignore')
                except:
                    return f"Binary file: {filename} (content extraction not supported)"
                    
        except Exception as e:
            return f"Error extracting content from {filename}: {str(e)}"
    
    def chunk_text(self, text: str, chunk_size: int = None, chunk_overlap: int = None) -> List[str]:
        """Split text into overlapping chunks"""
        if not text:
            return []
            
        chunk_size = chunk_size or self.chunk_size
        chunk_overlap = chunk_overlap or self.chunk_overlap
        
        chunks = []
        start = 0
        
        while start < len(text):
            end = start + chunk_size
            chunk = text[start:end]
            
            # Try to break at sentence boundary
            if end < len(text):
                # Look for sentence endings
                last_period = chunk.rfind('.')
                last_newline = chunk.rfind('\n')
                last_space = chunk.rfind(' ')
                
                # Use the best breaking point
                break_point = max(last_period, last_newline, last_space)
                if break_point > start + chunk_size // 2:  # Don't break too early
                    chunk = chunk[:break_point + 1]
                    end = start + len(chunk)
            
            chunks.append(chunk.strip())
            start = end - chunk_overlap
            
            if start >= len(text):
                break
        
        return [chunk for chunk in chunks if chunk.strip()]
    
    def compute_hash(self, content: str) -> str:
        """Compute hash of content for change detection"""
        return hashlib.md5(content.encode('utf-8')).hexdigest()
    
    def process_file(self, project_id: int, file_obj, filename: str, is_template: bool = False) -> Dict[str, Any]:
        """Process a single file: extract content, create embeddings, save to database"""
        if not self.is_available():
            return {"success": False, "error": "RAG service not available"}
        
        try:
            # Extract text content
            content = self.extract_text_from_file(file_obj, filename)
            if not content.strip():
                return {"success": False, "error": "No content extracted from file"}
            
            # Compute content hash
            content_hash = self.compute_hash(content)
            
            # Check if file already exists with same content
            existing_files = self.db_manager.get_project_data_files(project_id)
            for existing_file in existing_files:
                if (existing_file['filename'] == filename and 
                    existing_file['content_hash'] == content_hash):
                    return {
                        "success": True, 
                        "message": f"File {filename} already processed (no changes detected)",
                        "file_id": existing_file['id']
                    }
            
            # Save file to database
            file_path = f"project_{project_id}/{filename}"
            file_size = len(content.encode('utf-8'))
            file_type = filename.lower().split('.')[-1] if '.' in filename else 'unknown'
            
            file_id = self.db_manager.save_project_data_file(
                project_id=project_id,
                filename=filename,
                file_path=file_path,
                file_type=file_type,
                file_size=file_size,
                content=content,
                content_hash=content_hash,
                is_template=is_template
            )
            
            # Create text chunks
            chunks = self.chunk_text(content)
            
            # Generate embeddings for chunks
            embeddings_created = 0
            for chunk_index, chunk_text in enumerate(chunks):
                try:
                    # Generate embedding
                    embedding = self.model.encode(chunk_text).tolist()
                    
                    # Save embedding to database
                    metadata = {
                        "filename": filename,
                        "chunk_size": len(chunk_text),
                        "file_type": file_type,
                        "is_template": is_template
                    }
                    
                    self.db_manager.save_vector_embedding(
                        project_id=project_id,
                        file_id=file_id,
                        chunk_index=chunk_index,
                        chunk_text=chunk_text,
                        embedding_vector=embedding,
                        metadata=metadata
                    )
                    embeddings_created += 1
                    
                except Exception as e:
                    print(f"Error creating embedding for chunk {chunk_index}: {e}")
                    continue
            
            return {
                "success": True,
                "message": f"Successfully processed {filename}",
                "file_id": file_id,
                "chunks_created": len(chunks),
                "embeddings_created": embeddings_created,
                "content_preview": content[:200] + "..." if len(content) > 200 else content
            }
            
        except Exception as e:
            return {"success": False, "error": f"Error processing file {filename}: {str(e)}"}
    
    def process_folder(self, project_id: int, folder_path: str, is_template: bool = False) -> Dict[str, Any]:
        """Process all files in a folder recursively"""
        if not os.path.exists(folder_path):
            return {"success": False, "error": "Folder path does not exist"}
        
        results = []
        total_files = 0
        successful_files = 0
        
        # Supported file extensions
        supported_extensions = {
            '.txt', '.md', '.json', '.pdf', '.docx', '.doc', 
            '.xlsx', '.xls', '.pptx', '.ppt', '.csv'
        }
        
        for root, dirs, files in os.walk(folder_path):
            for file in files:
                file_path = os.path.join(root, file)
                file_extension = Path(file).suffix.lower()
                
                # Skip hidden files and unsupported types
                if file.startswith('.'):
                    continue
                
                total_files += 1
                
                try:
                    with open(file_path, 'rb') as file_obj:
                        result = self.process_file(project_id, file_obj, file, is_template)
                        results.append({
                            "filename": file,
                            "path": file_path,
                            "result": result
                        })
                        
                        if result["success"]:
                            successful_files += 1
                            
                except Exception as e:
                    results.append({
                        "filename": file,
                        "path": file_path,
                        "result": {"success": False, "error": str(e)}
                    })
        
        return {
            "success": True,
            "total_files": total_files,
            "successful_files": successful_files,
            "failed_files": total_files - successful_files,
            "results": results
        }
    
    def search_similar_content(self, project_id: int, query: str, top_k: int = 5) -> List[Dict[str, Any]]:
        """Search for similar content using vector similarity"""
        if not self.is_available():
            return []
        
        try:
            # Get all embeddings for the project
            embeddings_data = self.db_manager.get_vector_embeddings(project_id)
            
            if not embeddings_data:
                return []
            
            # Generate query embedding
            query_embedding = self.model.encode(query)
            
            # Compute similarities
            similarities = []
            for embedding_data in embeddings_data:
                stored_embedding = np.array(embedding_data['embedding_vector'])
                similarity = cosine_similarity(
                    query_embedding.reshape(1, -1), 
                    stored_embedding.reshape(1, -1)
                )[0][0]
                
                similarities.append({
                    **embedding_data,
                    'similarity': float(similarity)
                })
            
            # Sort by similarity and return top_k
            similarities.sort(key=lambda x: x['similarity'], reverse=True)
            return similarities[:top_k]
            
        except Exception as e:
            print(f"Error in similarity search: {e}")
            return []
    
    def get_context_for_query(self, project_id: int, query: str, max_context_length: int = 3000) -> str:
        """Get relevant context for a query using RAG"""
        similar_chunks = self.search_similar_content(project_id, query, top_k=5)
        
        if not similar_chunks:
            return ""
        
        context_parts = []
        current_length = 0
        
        for chunk_data in similar_chunks:
            chunk_text = chunk_data['chunk_text']
            filename = chunk_data['filename']
            similarity = chunk_data['similarity']
            
            # Add metadata
            context_part = f"[From {filename} (similarity: {similarity:.3f})]:\n{chunk_text}\n"
            
            if current_length + len(context_part) > max_context_length:
                break
                
            context_parts.append(context_part)
            current_length += len(context_part)
        
        return "\n---\n".join(context_parts)
    
    def get_project_data_summary(self, project_id: int) -> Dict[str, Any]:
        """Get summary of all project data"""
        files = self.db_manager.get_project_data_files(project_id)
        embeddings = self.db_manager.get_vector_embeddings(project_id)
        
        summary = {
            "total_files": len(files),
            "total_chunks": len(embeddings),
            "file_types": {},
            "template_files": 0,
            "data_files": 0
        }
        
        for file in files:
            file_type = file['file_type']
            summary["file_types"][file_type] = summary["file_types"].get(file_type, 0) + 1
            
            if file['is_template']:
                summary["template_files"] += 1
            else:
                summary["data_files"] += 1
        
        return summary
    
    def add_text_content(self, project_id: int, content: str, filename: str, file_type: str, is_template: bool = False) -> Dict:
        """Add text content directly to RAG system (for templates and other text content)"""
        if not self.is_available():
            return {"success": False, "error": "RAG service not available"}
        
        try:
            # Store file metadata
            file_data = {
                'project_id': project_id,
                'filename': filename,
                'file_type': file_type,
                'file_size': len(content.encode('utf-8')),
                'content': content,
                'is_template': is_template
            }
            
            file_id = self.db_manager.save_project_data_file(file_data)
            
            # Process content for embeddings
            chunks = self.chunk_text(content)
            embeddings_saved = 0
            
            for i, chunk in enumerate(chunks):
                try:
                    # Generate embedding
                    embedding = self.model.encode(chunk)
                    
                    # Save to database
                    embedding_data = {
                        'project_id': project_id,
                        'file_id': file_id,
                        'chunk_text': chunk,
                        'chunk_index': i,
                        'embedding': embedding.tolist()
                    }
                    
                    self.db_manager.save_vector_embedding(embedding_data)
                    embeddings_saved += 1
                    
                except Exception as e:
                    print(f"Warning: Could not process chunk {i} for {filename}: {e}")
                    continue
            
            return {
                "success": True,
                "file_id": file_id,
                "chunks_processed": embeddings_saved,
                "total_chunks": len(chunks)
            }
            
        except Exception as e:
            return {"success": False, "error": str(e)}
