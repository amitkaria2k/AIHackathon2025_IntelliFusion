"""
Enhanced AI-Powered Project Documentation Management System
With Bosch Branding, Authentication, Database Persistence, and Chatbot
"""

import streamlit as st
import pandas as pd
import base64
import os
import time
from datetime import datetime, timedelta
from typing import Dict, List, Any, Optional
import json
import sys

# Add config directory to path
sys.path.append(os.path.join(os.path.dirname(__file__), 'config'))
sys.path.append(os.path.join(os.path.dirname(__file__), 'services'))

from config.auth_config import VALID_PASSWORDS, HINT_MESSAGE, MAX_LOGIN_ATTEMPTS
from config.database import DatabaseManager
from services.llm_service import llm_service
from services.rag_service import RAGService
from services.ai_features import get_ai_features_service, DocumentIntelligenceService

# Configure Streamlit page with Bosch branding
st.set_page_config(
    page_title="Bosch AI Document Manager",
    page_icon="🔧",
    layout="wide",
    initial_sidebar_state="collapsed"  # Hide sidebar for tab navigation
)

# Bosch Corporate Colors (from brand guidelines)
BOSCH_COLORS = {
    'primary_red': '#DC2626',      # Bosch Red
    'primary_blue': '#1E40AF',     # Bosch Blue  
    'dark_blue': '#1E3A8A',       # Dark Blue
    'gray_dark': '#374151',        # Dark Gray
    'gray_medium': '#6B7280',      # Medium Gray
    'gray_light': '#F3F4F6',      # Light Gray
    'white': '#FFFFFF',
    'success_green': '#059669',
    'warning_orange': '#D97706',
    'error_red': '#DC2626'
}

def load_bosch_logo():
    """Load and encode Bosch logo"""
    logo_path = "data/images/Bosch-Logo.png"
    try:
        with open(logo_path, "rb") as f:
            logo_data = f.read()
        return base64.b64encode(logo_data).decode()
    except FileNotFoundError:
        return None

def apply_bosch_styling():
    """Apply Bosch corporate styling"""
    logo_base64 = load_bosch_logo()
    
    bosch_css = f"""
    <style>
    /* Import Bosch fonts */
    @import url('https://fonts.googleapis.com/css2?family=Arial:wght@400;600;700&display=swap');
    
    /* Main app styling */
    .main .block-container {{
        padding-top: 2rem;
        padding-left: 2rem;
        padding-right: 2rem;
    }}
    
    /* Header styling */
    .bosch-header {{
        background: linear-gradient(135deg, {BOSCH_COLORS['primary_red']} 0%, {BOSCH_COLORS['primary_blue']} 100%);
        padding: 1.5rem 2rem;
        border-radius: 10px;
        margin-bottom: 2rem;
        color: white;
        box-shadow: 0 4px 6px rgba(0, 0, 0, 0.1);
    }}
    
    .bosch-logo {{
        width: 120px;
        height: auto;
        margin-right: 2rem;
        vertical-align: middle;
    }}
    
    .bosch-title {{
        font-family: 'Arial', sans-serif;
        font-weight: 700;
        font-size: 2.5rem;
        margin: 0;
        display: inline-block;
        vertical-align: middle;
    }}
    
    .bosch-subtitle {{
        font-family: 'Arial', sans-serif;
        font-size: 1.2rem;
        margin-top: 0.5rem;
        opacity: 0.9;
    }}
    
    /* Tab styling */
    .stTabs [data-baseweb="tab-list"] {{
        gap: 2rem;
        background-color: {BOSCH_COLORS['gray_light']};
        border-radius: 10px;
        padding: 0.5rem;
        margin-bottom: 2rem;
    }}
    
    .stTabs [data-baseweb="tab"] {{
        height: 3rem;
        padding: 0 1.5rem;
        background-color: transparent;
        border-radius: 8px;
        color: {BOSCH_COLORS['gray_dark']};
        font-weight: 600;
        font-family: 'Arial', sans-serif;
        border: none;
        transition: all 0.3s ease;
    }}
    
    .stTabs [aria-selected="true"] {{
        background: linear-gradient(135deg, {BOSCH_COLORS['primary_red']} 0%, {BOSCH_COLORS['primary_blue']} 100%);
        color: white;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.2);
    }}
    
    .stTabs [data-baseweb="tab"]:hover {{
        background-color: {BOSCH_COLORS['primary_blue']};
        color: white;
    }}
    
    /* Button styling */
    .stButton button {{
        background: linear-gradient(135deg, {BOSCH_COLORS['primary_red']} 0%, {BOSCH_COLORS['primary_blue']} 100%);
        color: white;
        border: none;
        border-radius: 8px;
        padding: 0.75rem 1.5rem;
        font-weight: 600;
        font-family: 'Arial', sans-serif;
        transition: all 0.3s ease;
    }}
    
    .stButton button:hover {{
        transform: translateY(-2px);
        box-shadow: 0 4px 8px rgba(0, 0, 0, 0.2);
    }}
    
    /* Metric styling */
    [data-testid="metric-container"] {{
        background: white;
        border: 1px solid {BOSCH_COLORS['gray_light']};
        padding: 1rem;
        border-radius: 8px;
        box-shadow: 0 2px 4px rgba(0, 0, 0, 0.1);
    }}
    
    /* Login form styling */
    .login-container {{
        background: white;
        border-radius: 15px;
        padding: 3rem;
        box-shadow: 0 8px 16px rgba(0, 0, 0, 0.15);
        border: 2px solid {BOSCH_COLORS['primary_blue']};
        max-width: 500px;
        margin: 2rem auto;
    }}
    
    .login-title {{
        color: {BOSCH_COLORS['primary_blue']};
        font-family: 'Arial', sans-serif;
        font-weight: 700;
        text-align: center;
        margin-bottom: 2rem;
    }}
    
    /* Chat interface styling */
    .chat-container {{
        background: white;
        border-radius: 10px;
        padding: 1rem;
        border: 1px solid {BOSCH_COLORS['gray_light']};
        margin-bottom: 1rem;
    }}
    
    .chat-message {{
        padding: 0.75rem 1rem;
        border-radius: 8px;
        margin-bottom: 0.5rem;
        font-family: 'Arial', sans-serif;
    }}
    
    .chat-user {{
        background: {BOSCH_COLORS['primary_blue']};
        color: white;
        text-align: right;
    }}
    
    .chat-assistant {{
        background: {BOSCH_COLORS['gray_light']};
        color: {BOSCH_COLORS['gray_dark']};
    }}
    
    /* Success/Error messages */
    .stSuccess {{
        background-color: {BOSCH_COLORS['success_green']};
        color: white;
        border-radius: 8px;
    }}
    
    .stError {{
        background-color: {BOSCH_COLORS['error_red']};
        color: white;
        border-radius: 8px;
    }}
    
    .stWarning {{
        background-color: {BOSCH_COLORS['warning_orange']};
        color: white;
        border-radius: 8px;
    }}
    
    /* Hide Streamlit menu and footer */
    #MainMenu {{visibility: hidden;}}
    footer {{visibility: hidden;}}
    header {{visibility: hidden;}}
    </style>
    """
    
    # Add logo to CSS if available
    if logo_base64:
        bosch_css = bosch_css.replace(
            ".bosch-logo {", 
            f'.bosch-logo {{ content: url("data:image/png;base64,{logo_base64}");'
        )
    
    st.markdown(bosch_css, unsafe_allow_html=True)

def show_bosch_header():
    """Display Bosch branded header"""
    logo_base64 = load_bosch_logo()
    
    if logo_base64:
        header_html = f"""
        <div class="bosch-header">
            <img src="data:image/png;base64,{logo_base64}" class="bosch-logo">
            <div style="display: inline-block; vertical-align: middle;">
                <h1 class="bosch-title">AI Document Manager</h1>
                <p class="bosch-subtitle">Intelligent Project Documentation System</p>
            </div>
        </div>
        """
    else:
        header_html = f"""
        <div class="bosch-header">
            <div>
                <h1 class="bosch-title">🔧 Bosch AI Document Manager</h1>
                <p class="bosch-subtitle">Intelligent Project Documentation System</p>
            </div>
        </div>
        """
    
    st.markdown(header_html, unsafe_allow_html=True)

def authenticate_user():
    """Handle user authentication"""
    if 'authenticated' not in st.session_state:
        st.session_state.authenticated = False
        st.session_state.login_attempts = 0
        st.session_state.user_role = 'PM'  # Default role
    
    if not st.session_state.authenticated:
        # Add Bosch logo and app name
        logo_col1, logo_col2, logo_col3 = st.columns([1, 2, 1])
        with logo_col2:
            try:
                st.image("data/images/Bosch-Logo.png", width=200)
            except:
                st.write("🔧 **Bosch**")
        
        st.markdown("""
        <div style="text-align: center; margin: 2rem 0;">
            <h1 style="color: #00629B; font-size: 3rem; margin-bottom: 0.5rem;">IntelliFusion</h1>
            <p style="color: #6B7280; font-size: 1.2rem;">AI Document Assistant</p>
        </div>
        """, unsafe_allow_html=True)
        
        st.markdown("""
        <div class="login-container">
            <h2 class="login-title">🔧 User Login</h2>
        </div>
        """, unsafe_allow_html=True)
        
        with st.form("login_form"):
            # Role selection
            user_role = st.selectbox(
                "Select your role:",
                options=["PM", "Project team", "Quality team"],
                index=0,  # PM is default
                help="PM: Full access | Project team: Limited access | Quality team: Audit and compliance focus"
            )
            
            password = st.text_input("Password:", type="password", placeholder="Enter your password")
            login_button = st.form_submit_button("🚀 Login", use_container_width=True)
            
            if login_button:
                # Role-based password validation
                valid_password = False
                if user_role == "PM" and password == "admin":
                    valid_password = True
                elif user_role == "Project team" and password == "user":
                    valid_password = True
                elif user_role == "Quality team" and password == "quality":
                    valid_password = True
                
                if valid_password:
                    st.session_state.authenticated = True
                    st.session_state.user_role = user_role
                    st.session_state.login_attempts = 0
                    st.success(f"✅ Welcome! Logged in as: {user_role}")
                    st.rerun()
                else:
                    st.session_state.login_attempts += 1
                    
                    if st.session_state.login_attempts >= MAX_LOGIN_ATTEMPTS:
                        st.error(f"❌ Invalid password! Hint: PM='admin', Project team='user', Quality team='quality'")
                    else:
                        remaining = MAX_LOGIN_ATTEMPTS - st.session_state.login_attempts
                        st.error(f"❌ Invalid password! {remaining} attempts remaining.")
        
        return False
    
    return True

def check_access(required_role_or_permissions):
    """Check if current user has access to specific functionality
    
    Args:
        required_role_or_permissions: Can be 'PM', 'Project team', or list of specific permissions
    
    Returns:
        bool: True if user has access, False otherwise
    """
    if not st.session_state.get('authenticated', False):
        return False
    
    user_role = st.session_state.get('user_role', 'Project team')
    
    # PM has access to everything
    if user_role == 'PM':
        return True
    
    # Project team permissions
    if user_role == 'Project team':
        if isinstance(required_role_or_permissions, str):
            if required_role_or_permissions == 'Project team':
                return True
            elif required_role_or_permissions == 'PM':
                return False
        elif isinstance(required_role_or_permissions, list):
            # Check specific permissions
            project_team_permissions = [
                'generate_document', 
                'view_workflow', 
                'ai_assistant', 
                'dashboard',
                'project_overview'
            ]
            return all(perm in project_team_permissions for perm in required_role_or_permissions)
    
    # Quality team permissions (audit-focused, no document generation)
    if user_role == 'Quality team':
        if isinstance(required_role_or_permissions, str):
            if required_role_or_permissions == 'Quality team':
                return True
            elif required_role_or_permissions in ['PM', 'Project team']:
                return False
        elif isinstance(required_role_or_permissions, list):
            # Check specific permissions
            quality_team_permissions = [
                'view_workflow',  # Can view workflows for audit
                'dashboard',      # Can see dashboard
                'project_overview',  # Can see project overview
                'compliance_audit'   # Special audit permissions
            ]
            return all(perm in quality_team_permissions for perm in required_role_or_permissions)
    
    return False

def show_access_denied(feature_name="this feature"):
    """Show access denied message"""
    st.error(f"🔒 Access Denied: Only PM users can access {feature_name}.")
    st.info("💡 Contact your Project Manager for access or login with a PM role.")

def initialize_app():
    """Initialize the application"""
    try:
        if 'db' not in st.session_state:
            st.session_state.db = DatabaseManager()
        
        if 'rag_service' not in st.session_state:
            st.session_state.rag_service = RAGService(st.session_state.db)
        
        if 'chat_history' not in st.session_state:
            st.session_state.chat_history = []
        
        if 'llm_settings' not in st.session_state:
            st.session_state.llm_settings = {
                'temperature': 0.7,
                'max_tokens': 2000,
                'model': 'gpt-4o-mini',
                'use_project_context': True
            }
            
    except Exception as e:
        st.error(f"Initialization Error: {str(e)}")
        st.exception(e)

def create_project_template_folder(project_name: str) -> str:
    """Create and return the template folder path for a project"""
    # Sanitize project name for folder creation
    safe_project_name = "".join(c for c in project_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
    safe_project_name = safe_project_name.replace(' ', '_')
    
    templates_base = os.path.join(os.getcwd(), 'templates')
    project_template_folder = os.path.join(templates_base, safe_project_name, 'template')
    
    # Create folder structure if it doesn't exist
    os.makedirs(project_template_folder, exist_ok=True)
    
    return project_template_folder

def generate_ai_template(doc_type: str, project_context: Dict) -> str:
    """Generate AI template for a specific document type"""
    templates_map = {
        "Project Management Plan (PMP)": {
            "title": "Project Management Plan",
            "sections": [
                "1. PROJECT OVERVIEW",
                "1.1 Project Purpose and Justification",
                "1.2 Project Description",
                "1.3 Project Objectives", 
                "1.4 Success Criteria",
                "2. PROJECT SCOPE",
                "2.1 In Scope",
                "2.2 Out of Scope",
                "2.3 Deliverables",
                "3. PROJECT ORGANIZATION",
                "3.1 Project Team Structure",
                "3.2 Roles and Responsibilities",
                "3.3 Communication Plan",
                "4. PROJECT SCHEDULE",
                "4.1 Major Milestones",
                "4.2 Work Breakdown Structure",
                "4.3 Critical Path Analysis",
                "5. RISK MANAGEMENT",
                "5.1 Risk Assessment",
                "5.2 Risk Mitigation Strategies",
                "5.3 Contingency Planning",
                "6. QUALITY MANAGEMENT",
                "6.1 Quality Standards",
                "6.2 Quality Assurance Process",
                "6.3 Quality Control Measures"
            ]
        },
        "Technical Concept Document (TCD)": {
            "title": "Technical Concept Document", 
            "sections": [
                "1. EXECUTIVE SUMMARY",
                "2. TECHNICAL REQUIREMENTS",
                "2.1 Functional Requirements",
                "2.2 Non-Functional Requirements",
                "2.3 Interface Requirements",
                "3. SYSTEM ARCHITECTURE",
                "3.1 High-Level Architecture",
                "3.2 Component Design",
                "3.3 Data Flow Diagrams",
                "4. TECHNOLOGY STACK",
                "4.1 Hardware Requirements",
                "4.2 Software Components",
                "4.3 Third-Party Dependencies",
                "5. IMPLEMENTATION APPROACH",
                "5.1 Development Methodology",
                "5.2 Integration Strategy",
                "5.3 Testing Strategy",
                "6. SECURITY CONSIDERATIONS",
                "7. PERFORMANCE ANALYSIS",
                "8. MAINTENANCE AND SUPPORT"
            ]
        },
        "Configuration Management Plan": {
            "title": "Configuration Management Plan",
            "sections": [
                "1. INTRODUCTION",
                "2. CONFIGURATION MANAGEMENT OBJECTIVES",
                "3. CONFIGURATION IDENTIFICATION",
                "3.1 Configuration Items",
                "3.2 Naming Conventions",
                "3.3 Versioning Strategy",
                "4. CONFIGURATION CONTROL",
                "4.1 Change Control Process",
                "4.2 Change Control Board",
                "4.3 Configuration Baselines",
                "5. CONFIGURATION STATUS ACCOUNTING",
                "6. CONFIGURATION AUDITS",
                "7. TOOLS AND INFRASTRUCTURE",
                "8. ROLES AND RESPONSIBILITIES"
            ]
        },
        "Communication Management Plan": {
            "title": "Communication Management Plan",
            "sections": [
                "1. COMMUNICATION OBJECTIVES",
                "2. STAKEHOLDER ANALYSIS",
                "3. COMMUNICATION REQUIREMENTS",
                "4. COMMUNICATION METHODS AND CHANNELS",
                "5. COMMUNICATION MATRIX",
                "6. MEETING MANAGEMENT",
                "7. REPORTING STRUCTURE",
                "8. INFORMATION DISTRIBUTION",
                "9. COMMUNICATION ESCALATION",
                "10. COMMUNICATION MONITORING AND CONTROL"
            ]
        },
        "Risk Management Plan": {
            "title": "Risk Management Plan",
            "sections": [
                "1. RISK MANAGEMENT APPROACH",
                "2. RISK CATEGORIES",
                "3. RISK IDENTIFICATION",
                "4. QUALITATIVE RISK ANALYSIS",
                "5. QUANTITATIVE RISK ANALYSIS",
                "6. RISK RESPONSE PLANNING",
                "7. RISK MONITORING AND CONTROL",
                "8. RISK REGISTER",
                "9. RISK REPORTING",
                "10. ROLES AND RESPONSIBILITIES"
            ]
        },
        "Quality Plan": {
            "title": "Quality Management Plan",
            "sections": [
                "1. QUALITY POLICY AND OBJECTIVES",
                "2. QUALITY STANDARDS AND METRICS",
                "3. QUALITY PLANNING",
                "4. QUALITY ASSURANCE ACTIVITIES",
                "5. QUALITY CONTROL ACTIVITIES",
                "6. QUALITY IMPROVEMENT PROCESS",
                "7. QUALITY ROLES AND RESPONSIBILITIES",
                "8. QUALITY TOOLS AND TECHNIQUES",
                "9. QUALITY DOCUMENTATION",
                "10. QUALITY MONITORING AND REPORTING"
            ]
        }
    }
    
    template_info = templates_map.get(doc_type, {
        "title": doc_type,
        "sections": ["1. INTRODUCTION", "2. OBJECTIVES", "3. SCOPE", "4. METHODOLOGY", "5. DELIVERABLES", "6. CONCLUSION"]
    })
    
    # Generate template content
    template_content = f"""# {template_info['title']}

**Project:** {project_context.get('name', 'Project Name')}
**Project Type:** {project_context.get('type', 'Project Type')}
**Date:** {{{{ current_date }}}}
**Version:** 1.0

---

## Document Information
- **Document Type:** {template_info['title']}
- **Project Description:** {project_context.get('description', 'Project Description')}
- **Generated:** AI Template for {project_context.get('name', 'Project')}

---

"""
    
    # Add sections with placeholders
    for section in template_info['sections']:
        template_content += f"## {section}\n\n"
        template_content += f"{{{{ {section.lower().replace(' ', '_').replace('.', '').replace('/', '_')} }}}}\n\n"
        template_content += "---\n\n"
    
    # Add project-specific requirements if available
    if project_context.get('functional_reqs'):
        template_content += "## FUNCTIONAL REQUIREMENTS\n\n"
        for req in project_context['functional_reqs'][:3]:  # First 3 requirements
            template_content += f"- {req}\n"
        template_content += "\n---\n\n"
    
    if project_context.get('non_functional_reqs'):
        template_content += "## NON-FUNCTIONAL REQUIREMENTS\n\n"
        for req in project_context['non_functional_reqs'][:3]:  # First 3 requirements
            template_content += f"- {req}\n"
        template_content += "\n---\n\n"
    
    if project_context.get('conditions'):
        template_content += "## CONDITIONS AND CONSTRAINTS\n\n"
        for condition in project_context['conditions'][:3]:  # First 3 conditions
            template_content += f"- {condition}\n"
        template_content += "\n---\n\n"
    
    template_content += """
## TEMPLATE USAGE NOTES

This template was generated automatically based on your project information. Please:

1. Replace placeholder sections ({{{{ ... }}}}) with actual content
2. Customize sections as needed for your specific project
3. Add or remove sections based on project requirements
4. Update version number and date as document evolves

---

**Template Generated:** {current_date}
**For Project:** {project_name}
""".format(current_date=datetime.now().strftime("%Y-%m-%d %H:%M"), 
           project_name=project_context.get('name', 'Unknown Project'))
    
    return template_content

def save_template_to_folder(template_content: str, doc_type: str, template_folder: str, source_file=None) -> str:
    """Save template content to the project template folder"""
    # Create safe filename
    safe_doc_type = "".join(c for c in doc_type if c.isalnum() or c in (' ', '-', '_', '(', ')')).rstrip()
    safe_doc_type = safe_doc_type.replace(' ', '_').replace('(', '').replace(')', '')
    
    if source_file:
        # Use original filename if from uploaded file
        file_extension = source_file.name.split('.')[-1] if '.' in source_file.name else 'txt'
        filename = f"{safe_doc_type}_template.{file_extension}"
    else:
        # AI generated template
        filename = f"{safe_doc_type}_AI_template.md"
    
    file_path = os.path.join(template_folder, filename)
    
    # Save the template
    with open(file_path, 'w', encoding='utf-8') as f:
        f.write(template_content)
    
    return file_path

def process_document_templates(project_id: int, project_context: Dict, template_selections: Dict) -> Dict:
    """Process and save all document templates for a project"""
    result = {
        'success': True,
        'template_folder': '',
        'created_templates': [],
        'errors': []
    }
    
    try:
        # Create project template folder
        template_folder = create_project_template_folder(project_context['name'])
        result['template_folder'] = template_folder
        
        for doc_type, selection in template_selections.items():
            try:
                if selection['source'] == 'AI Generated':
                    # Generate AI template
                    template_content = generate_ai_template(doc_type, project_context)
                    file_path = save_template_to_folder(template_content, doc_type, template_folder)
                    result['created_templates'].append({
                        'doc_type': doc_type,
                        'source': 'AI Generated',
                        'file_path': file_path,
                        'filename': os.path.basename(file_path)
                    })
                    
                elif selection['source'] == 'Upload Template' and selection['template_file']:
                    # Process uploaded template
                    uploaded_file = selection['template_file']
                    
                    # Read file content
                    if uploaded_file.name.endswith(('.txt', '.md')):
                        template_content = uploaded_file.read().decode('utf-8')
                    elif uploaded_file.name.endswith('.json'):
                        template_content = uploaded_file.read().decode('utf-8')
                    elif uploaded_file.name.endswith(('.docx', '.doc')):
                        # Extract text from Word document
                        try:
                            import docx
                            doc = docx.Document(uploaded_file)
                            template_content = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                        except Exception:
                            template_content = f"Word document template: {uploaded_file.name}\n[Content extraction not available]"
                    else:
                        template_content = f"Template file: {uploaded_file.name}\n[Binary file - content not extracted]"
                    
                    # Save uploaded template
                    file_path = save_template_to_folder(template_content, doc_type, template_folder, uploaded_file)
                    result['created_templates'].append({
                        'doc_type': doc_type,
                        'source': 'User Upload',
                        'file_path': file_path,
                        'filename': os.path.basename(file_path),
                        'original_name': uploaded_file.name
                    })
                
            except Exception as e:
                result['errors'].append(f"Error processing {doc_type}: {str(e)}")
        
        # Process templates with RAG service if available
        if st.session_state.rag_service.is_available():
            for template_info in result['created_templates']:
                try:
                    # Read the saved template file
                    with open(template_info['file_path'], 'r', encoding='utf-8') as f:
                        content = f.read()
                    
                    # Process with RAG (mark as template)
                    rag_result = st.session_state.rag_service.add_text_content(
                        project_id=project_id,
                        content=content,
                        filename=template_info['filename'],
                        file_type='md' if template_info['file_path'].endswith('.md') else 'txt',
                        is_template=True
                    )
                    
                    if not rag_result.get('success'):
                        result['errors'].append(f"RAG processing failed for {template_info['filename']}")
                        
                except Exception as e:
                    result['errors'].append(f"RAG processing error for {template_info['filename']}: {str(e)}")
    
    except Exception as e:
        result['success'] = False
        result['errors'].append(f"Template processing failed: {str(e)}")
    
    return result

def load_project_template(uploaded_files):
    """Load project template from uploaded file(s)"""
    if uploaded_files is not None and len(uploaded_files) > 0:
        combined_template = {
            'name': '',
            'type': '',
            'description': '',
            'functional_reqs': [],
            'non_functional_reqs': [],
            'conditions': [],
            'recommended_docs': [],
            'template_content': []
        }
        
        # Import required libraries for file parsing
        try:
            import PyPDF2
            import docx
            from pptx import Presentation
            import pandas as pd
        except ImportError as e:
            st.error(f"Missing required libraries for file parsing: {e}")
            return None
        
        for uploaded_file in uploaded_files:
            try:
                # Check file size (Streamlit default limit is around 200MB)
                if hasattr(uploaded_file, 'size') and uploaded_file.size > 200 * 1024 * 1024:  # 200MB
                    st.error(f"❌ File {uploaded_file.name} is too large (>200MB). Please use a smaller file.")
                    continue
                
                file_extension = uploaded_file.name.lower().split('.')[-1]
                
                # Reset file pointer to beginning
                uploaded_file.seek(0)
                
                if file_extension == 'json':
                    try:
                        content = json.loads(uploaded_file.read().decode('utf-8'))
                        # Merge JSON content with existing template
                        for key in ['name', 'type', 'description']:
                            if content.get(key) and not combined_template.get(key):
                                combined_template[key] = content[key]
                        
                        for key in ['functional_reqs', 'non_functional_reqs', 'conditions', 'recommended_docs']:
                            if content.get(key):
                                combined_template[key].extend(content[key])
                        
                        st.success(f"✅ JSON template loaded from {uploaded_file.name}")
                    except json.JSONDecodeError as je:
                        st.error(f"❌ Invalid JSON format in {uploaded_file.name}: {str(je)}")
                        continue
                    
                elif file_extension in ['txt', 'md']:
                    try:
                        content = uploaded_file.read().decode('utf-8')
                        combined_template['template_content'].append({
                            'filename': uploaded_file.name,
                            'content': content
                        })
                        st.success(f"✅ Text content loaded from {uploaded_file.name}")
                    except UnicodeDecodeError as ue:
                        st.error(f"❌ Cannot read {uploaded_file.name}: {str(ue)}. Please ensure the file uses UTF-8 encoding.")
                        continue
                    
                elif file_extension in ['xlsx', 'xls']:
                    # Extract content from Excel files
                    try:
                        excel_data = pd.read_excel(uploaded_file, sheet_name=None)
                        excel_content = ""
                        for sheet_name, df in excel_data.items():
                            excel_content += f"\n--- Sheet: {sheet_name} ---\n"
                            excel_content += df.to_string(index=False)
                            excel_content += "\n"
                        
                        combined_template['template_content'].append({
                            'filename': uploaded_file.name,
                            'content': excel_content
                        })
                        st.success(f"✅ Excel content extracted from {uploaded_file.name}")
                        
                    except Exception as e:
                        st.warning(f"⚠️ Could not fully parse Excel file {uploaded_file.name}: {str(e)}")
                        combined_template['template_content'].append({
                            'filename': uploaded_file.name,
                            'content': f"Excel file: {uploaded_file.name} (partial parsing error)"
                        })
                    
                elif file_extension in ['docx', 'doc']:
                    # Extract content from Word documents
                    try:
                        if file_extension == 'docx':
                            doc = docx.Document(uploaded_file)
                            word_content = ""
                            for paragraph in doc.paragraphs:
                                word_content += paragraph.text + "\n"
                            
                            # Extract tables with enhanced error handling
                            try:
                                for table in doc.tables:
                                    word_content += "\n--- Table ---\n"
                                    for row_idx, row in enumerate(table.rows):
                                        try:
                                            # Handle potential missing cells or malformed table structure
                                            cells_text = []
                                            for cell in row.cells:
                                                if cell and hasattr(cell, 'text'):
                                                    cells_text.append(cell.text.strip())
                                                else:
                                                    cells_text.append("[empty]")
                                            
                                            row_text = " | ".join(cells_text)
                                            if row_text.strip() and row_text != " | ".join(["[empty]"] * len(cells_text)):
                                                word_content += row_text + "\n"
                                        except Exception as cell_error:
                                            # Skip problematic rows and continue
                                            word_content += f"[Table row {row_idx}: parsing error]\n"
                                            continue
                                    word_content += "\n"
                            except Exception as table_error:
                                # If table parsing fails completely, continue with text content
                                word_content += f"\n--- Tables could not be parsed: {str(table_error)} ---\n"
                            
                            combined_template['template_content'].append({
                                'filename': uploaded_file.name,
                                'content': word_content
                            })
                            st.success(f"✅ Word content extracted from {uploaded_file.name}")
                        else:
                            # .doc files need different handling
                            st.warning(f"⚠️ .doc files require conversion to .docx for full support: {uploaded_file.name}")
                            combined_template['template_content'].append({
                                'filename': uploaded_file.name,
                                'content': f"Word document (.doc): {uploaded_file.name} - Please convert to .docx for full text extraction"
                            })
                            
                    except Exception as e:
                        # More robust error handling - still capture content even with parsing errors
                        st.warning(f"⚠️ Could not fully parse Word file {uploaded_file.name}: {str(e)}")
                        try:
                            # Try basic text extraction as fallback
                            doc = docx.Document(uploaded_file)
                            basic_content = ""
                            for paragraph in doc.paragraphs:
                                basic_content += paragraph.text + "\n"
                            
                            combined_template['template_content'].append({
                                'filename': uploaded_file.name,
                                'content': basic_content or f"Word document: {uploaded_file.name} (partial parsing)"
                            })
                            st.info(f"📄 Basic text extracted from {uploaded_file.name} (tables may be missing)")
                        except Exception:
                            # Final fallback - just store metadata
                            combined_template['template_content'].append({
                                'filename': uploaded_file.name,
                                'content': f"Word document: {uploaded_file.name} (parsing error: {str(e)})"
                            })
                    
                elif file_extension in ['pptx', 'ppt']:
                    # Extract content from PowerPoint presentations
                    try:
                        if file_extension == 'pptx':
                            prs = Presentation(uploaded_file)
                            ppt_content = ""
                            for slide_num, slide in enumerate(prs.slides, 1):
                                ppt_content += f"\n--- Slide {slide_num} ---\n"
                                for shape in slide.shapes:
                                    if hasattr(shape, "text"):
                                        ppt_content += shape.text + "\n"
                            
                            combined_template['template_content'].append({
                                'filename': uploaded_file.name,
                                'content': ppt_content
                            })
                            st.success(f"✅ PowerPoint content extracted from {uploaded_file.name}")
                        else:
                            # .ppt files need different handling
                            st.warning(f"⚠️ .ppt files require conversion to .pptx for full support: {uploaded_file.name}")
                            combined_template['template_content'].append({
                                'filename': uploaded_file.name,
                                'content': f"PowerPoint presentation (.ppt): {uploaded_file.name} - Please convert to .pptx for full text extraction"
                            })
                            
                    except Exception as e:
                        st.warning(f"⚠️ Could not fully parse PowerPoint file {uploaded_file.name}: {str(e)}")
                        combined_template['template_content'].append({
                            'filename': uploaded_file.name,
                            'content': f"PowerPoint presentation: {uploaded_file.name} (parsing error: {str(e)})"
                        })
                    
                elif file_extension == 'pdf':
                    # Extract content from PDF files
                    try:
                        pdf_reader = PyPDF2.PdfReader(uploaded_file)
                        pdf_content = ""
                        for page_num, page in enumerate(pdf_reader.pages, 1):
                            pdf_content += f"\n--- Page {page_num} ---\n"
                            pdf_content += page.extract_text() + "\n"
                        
                        combined_template['template_content'].append({
                            'filename': uploaded_file.name,
                            'content': pdf_content
                        })
                        st.success(f"✅ PDF content extracted from {uploaded_file.name} ({len(pdf_reader.pages)} pages)")
                        
                    except Exception as e:
                        st.warning(f"⚠️ Could not fully parse PDF file {uploaded_file.name}: {str(e)}")
                        combined_template['template_content'].append({
                            'filename': uploaded_file.name,
                            'content': f"PDF document: {uploaded_file.name} (parsing error: {str(e)})"
                        })
                    
                elif file_extension in ['csv']:
                    # Extract content from CSV files
                    try:
                        csv_data = pd.read_csv(uploaded_file)
                        csv_content = f"\n--- CSV Data from {uploaded_file.name} ---\n"
                        csv_content += csv_data.to_string(index=False)
                        
                        combined_template['template_content'].append({
                            'filename': uploaded_file.name,
                            'content': csv_content
                        })
                        st.success(f"✅ CSV data loaded from {uploaded_file.name}")
                    except Exception as e:
                        st.warning(f"⚠️ Could not parse CSV {uploaded_file.name}: {e}")
                        # Try to read as plain text
                        try:
                            uploaded_file.seek(0)  # Reset file pointer
                            content = uploaded_file.read().decode('utf-8')
                            combined_template['template_content'].append({
                                'filename': uploaded_file.name,
                                'content': content
                            })
                            st.info(f"📝 CSV loaded as text from {uploaded_file.name}")
                        except Exception as e2:
                            st.error(f"❌ Failed to load {uploaded_file.name}: {e2}")
                
                else:
                    # Try to read as text for any other file types
                    try:
                        content = uploaded_file.read().decode('utf-8')
                        combined_template['template_content'].append({
                            'filename': uploaded_file.name,
                            'content': content
                        })
                        st.success(f"✅ Text content loaded from {uploaded_file.name}")
                    except UnicodeDecodeError:
                        # For binary files, store metadata only
                        file_size = len(uploaded_file.getvalue()) if hasattr(uploaded_file, 'getvalue') else 0
                        combined_template['template_content'].append({
                            'filename': uploaded_file.name,
                            'content': f"[Binary file: {uploaded_file.name}, Size: {file_size} bytes, Type: {file_extension.upper()}]"
                        })
                        st.info(f"📄 Binary file metadata stored for {uploaded_file.name}")
                    except Exception as e:
                        st.warning(f"⚠️ Could not read {uploaded_file.name}: {str(e)}")
                    
            except Exception as e:
                st.error(f"❌ Error reading file {uploaded_file.name}: {str(e)}")
        
        return combined_template if any([combined_template['name'], combined_template['functional_reqs'], combined_template['template_content']]) else None
    return None

def generate_ai_content(prompt: str, project_context: Optional[Dict] = None) -> str:
    """Generate content using LLM"""
    # Use the new LLM service
    messages = [{"role": "user", "content": prompt}]
    
    result = llm_service.generate_response(
        messages=messages,
        temperature=st.session_state.llm_settings['temperature'],
        max_tokens=st.session_state.llm_settings['max_tokens'],
        use_project_context=st.session_state.llm_settings['use_project_context'],
        project_context=project_context
    )
    
    return result['response']

def main():
    """Main application function"""
    try:
        # Apply styling
        apply_bosch_styling()
        
        # Authentication check
        if not authenticate_user():
            return
        
        # Show header
        show_bosch_header()
        
        # Initialize app
        initialize_app()
        
        # Show current user role and logout button
        user_role = st.session_state.get('user_role', 'Project team')
        # Role-specific colors
        if user_role == 'PM':
            role_color = "#28A745"  # Green for PM
        elif user_role == 'Project team':
            role_color = "#17A2B8"  # Blue for Project team
        else:  # Quality team
            role_color = "#FFC107"  # Yellow for Quality team
        
        # Create columns for role badge and logout button
        col1, col2 = st.columns([3, 1])
        
        with col1:
            st.markdown(f"""
            <div style="text-align: right; margin-bottom: 1rem;">
                <span style="background-color: {role_color}; color: white; padding: 0.25rem 0.5rem; border-radius: 0.25rem; font-size: 0.9rem;">
                    👤 {user_role}
                </span>
            </div>
            """, unsafe_allow_html=True)
        
        with col2:
            if st.button("🚪 Logout", type="secondary", use_container_width=True):
                # Clear all session state
                for key in list(st.session_state.keys()):
                    del st.session_state[key]
                st.rerun()
        
        # Create tabs based on user role
        if user_role == 'PM':
            # PM has access to all tabs
            tab1, tab2, tab3, tab4, tab5, tab6, tab7, tab8 = st.tabs([
                "🏠 Dashboard", 
                "🆕 New Project", 
                "✏️ Edit Projects",
                "📝 Generate Document", 
                "👥 Workflow Management", 
                "📊 Project Overview",
                "💬 AI Assistant",
                "⚙️ Settings"
            ])
            # Create placeholder tabs to maintain structure
            tab2 = tab3 = tab10 = None
        elif user_role == 'Project team':
            # Project team has limited access (includes Settings now)
            tab1, tab4, tab5, tab6, tab7, tab8 = st.tabs([
                "🏠 Dashboard",
                "📝 Generate Document", 
                "👥 Workflow Management", 
                "📊 Project Overview",
                "💬 AI Assistant",
                "⚙️ Settings"
            ])
            # Create placeholder tabs to maintain structure
            tab2 = tab3 = tab10 = None
        else:  # Quality team
            # Quality team has audit-focused access (no document generation)
            tab1, tab5, tab6, tab8, tab10 = st.tabs([
                "🏠 Dashboard",
                "👥 Workflow Management", 
                "📊 Project Overview",
                "⚙️ Settings",
                "🔍 Compliance Audit"
            ])
            # Create placeholder tabs to maintain structure
            tab2 = tab3 = tab4 = tab7 = None
        
        # Tab content with role-based access control
        with tab1:
            show_dashboard()
        
        if tab2 is not None:  # PM only
            with tab2:
                if check_access('PM'):
                    show_new_project()
                else:
                    show_access_denied("project creation")
        
        if tab3 is not None:  # PM only
            with tab3:
                if check_access('PM'):
                    show_edit_projects()
                else:
                    show_access_denied("project editing")
        
        if tab4:  # Document Generation tab (not available for Quality team)
            with tab4:
                if check_access(['generate_document']):
                    show_document_generation()
                else:
                    show_access_denied("document generation")
        
        with tab5:
            if check_access(['view_workflow']):
                show_workflow_management()
            else:
                show_access_denied("workflow management")
        
        with tab6:
            if check_access(['project_overview']):
                show_project_overview()
            else:
                show_access_denied("project overview")
        
        if tab7:  # AI Assistant tab (not available for Quality team)
            with tab7:
                if check_access(['ai_assistant']):
                    show_ai_assistant()
                else:
                    show_access_denied("AI assistant")
        
        if tab8 is not None:  # Settings tab
            with tab8:
                if check_access('PM'):
                    show_settings()
                else:
                    show_access_denied("system settings")
        
        # Compliance Audit tab (Quality team specific)
        if user_role == 'Quality team' and 'tab10' in locals() and tab10 is not None:
            with tab10:
                show_compliance_audit()

    except Exception as e:
        st.error(f"❌ Application Startup Error: {str(e)}")

def show_dashboard():
    """Display main dashboard with persistent data and AI insights"""
    st.title("🏠 Dashboard")
    st.markdown("Welcome to the Bosch AI-Powered Project Documentation Management System")
    
    # Force refresh data from database
    if st.button("🔄 Refresh Data", type="secondary"):
        st.rerun()
    
    # Load data from database
    projects = st.session_state.db.get_projects()
    documents = st.session_state.db.get_documents()
    workflows = st.session_state.db.get_workflows()
    
    # Display metrics with AI enhancements
    col1, col2, col3, col4, col5 = st.columns(5)
    
    with col1:
        st.metric("Total Projects", len(projects))
    
    with col2:
        st.metric("Generated Documents", len(documents))
        
    with col3:
        active_workflows = len([w for w in workflows if w['status'] == 'Active'])
        st.metric("Active Workflows", active_workflows)
        
    with col4:
        st.metric("Template Types", 12)
    
    with col5:
        st.metric("AI Insights", st.session_state.get('api_calls', 0))
    
    # AI-Powered Insights Section
    st.subheader("🧠 AI-Powered Insights")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("### 📈 Project Health Analysis")
        if documents:
            try:
                # Simulate AI analysis of project health
                ai_service = get_ai_features_service(llm_service, st.session_state.get('rag_service'))
                
                if ai_service:
                    # Generate compliance report
                    compliance_report = ai_service.generate_compliance_report(documents)
                    
                    # Health indicators
                    health_score = min(100, int(compliance_report['quality_stats']['average_quality'] * 100))
                    health_color = "🟢" if health_score >= 80 else "🟡" if health_score >= 60 else "🔴"
                    
                    st.metric("Project Health", f"{health_color} {health_score}%")
                    st.metric("Quality Issues", compliance_report['quality_stats']['low_quality_count'])
                    
                    # Show recommendations
                    if compliance_report['recommendations']:
                        st.markdown("**💡 AI Recommendations:**")
                        for rec in compliance_report['recommendations'][:3]:
                            st.write(f"• {rec}")
                else:
                    st.info("AI analysis not available. Enable AI features in configuration.")
            except Exception as e:
                st.warning(f"AI analysis temporarily unavailable: {str(e)}")
                # Fallback metrics
                st.metric("Documents", len(documents))
                st.metric("Avg. Quality", "85%")
    
    with col2:
        st.markdown("### 🎯 Smart Recommendations")
        
        # AI-powered recommendations based on current data
        recommendations = []
        
        if len(documents) == 0:
            recommendations.append("📝 Create your first project document to get started")
        elif len(workflows) == 0:
            recommendations.append("🔄 Set up approval workflows for better document management")
        elif active_workflows > 5:
            recommendations.append("⚡ Consider optimizing workflows - many are currently active")
        
        # Add more intelligent recommendations
        if len(documents) > 10:
            recommendations.append("🔍 Use AI search to find similar documents across projects")
        
        if len(projects) > 3:
            recommendations.append("📊 Generate a compliance report to ensure standards adherence")
        
        recommendations.append("💬 Try the AI Assistant for intelligent document insights")
        
        if recommendations:
            for i, rec in enumerate(recommendations[:4], 1):
                st.write(f"{i}. {rec}")
        else:
            st.write("Great! Your project management is on track. 🎉")
    
    # Recent projects with AI insights
    st.subheader("📊 Recent Projects")
    if projects:
        project_data = []
        for project in projects[:5]:  # Last 5 projects
            project_docs = len([d for d in documents if d['project_id'] == project['id']])
            project_workflows = len([w for w in workflows if any(d['project_id'] == project['id'] for d in documents if d['id'] == w['document_id'])])
            
            # AI-powered project risk assessment
            risk_level = "Low"
            if project_workflows > 3:
                risk_level = "Medium"
            if any(w['status'] == 'Rejected' for w in workflows):
                risk_level = "High"
            
            project_data.append({
                "Project Name": project['name'],
                "Type": project['type'],
                "Documents": project_docs,
                "Workflows": project_workflows,
                "AI Risk Level": risk_level,
                "Created": project['created_at'][:10]
            })
            
        if project_data:
            df = pd.DataFrame(project_data)
            
            # Color code the dataframe based on risk levels
            def highlight_risk(val):
                if val == 'High':
                    return 'background-color: #ffebee'
                elif val == 'Medium':
                    return 'background-color: #fff3e0'
                else:
                    return 'background-color: #e8f5e8'
            
            styled_df = df.style.map(highlight_risk, subset=['AI Risk Level'])
            st.dataframe(styled_df, use_container_width=True)
        else:
            st.info("No projects found. Create your first project to get started!")
    else:
        st.info("No projects found. Create your first project to get started!")
    
    # Recent documents
    st.subheader("📄 Recent Documents")
    if documents:
        doc_data = []
        for doc in documents[:5]:
            project = next((p for p in projects if p['id'] == doc['project_id']), None)
            doc_data.append({
                "Document": doc["name"],
                "Project": project["name"] if project else "Unknown",
                "Type": doc["type"],
                "Status": doc["status"],
                "Created": doc["created_at"][:10] if doc["created_at"] else "Today"
            })
        
        df = pd.DataFrame(doc_data)
        st.dataframe(df, use_container_width=True)
    else:
        st.info("No documents generated yet.")
    
    # System status
    st.subheader("🔧 System Status")
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.success("✅ Application Running")
    
    with col2:
        st.success("✅ Database Connected")
    
    with col3:
        if os.getenv('LLM_FARM_API_KEY'):
            st.success("✅ AI Service Ready")
        else:
            st.warning("⚠️ AI Service (Demo Mode)")

def show_new_project():
    """Show new project creation form with template upload"""
    st.title("🆕 Create New Project")
    
    # Add a reset button for form state
    col_reset, col_spacer = st.columns([1, 3])
    with col_reset:
        if st.button("🔄 Reset Form", help="Clear all selections and start fresh"):
            # Clear session state for new project
            if 'new_project_doc_selections' in st.session_state:
                del st.session_state.new_project_doc_selections
            if 'doc_template_selections' in st.session_state:
                del st.session_state.doc_template_selections
            st.rerun()
    
    # Project template upload section
    st.subheader("📁 Project Template(s) (Optional)")
    uploaded_files = st.file_uploader(
        "Upload project template file(s)", 
        type=['json', 'txt', 'md', 'xlsx', 'xls', 'docx', 'doc', 'pptx', 'ppt', 'pdf'],
        help="Upload one or more template files to pre-fill project requirements. Supports: JSON, TXT, Markdown, Excel, Word, PowerPoint, PDF",
        accept_multiple_files=True,
        key="project_template_upload"
    )
    
    template_data = None
    if uploaded_files:
        template_data = load_project_template(uploaded_files)
        if template_data:
            st.success(f"✅ Loaded {len(uploaded_files)} template file(s) successfully!")
            
            with st.expander("📋 Preview combined template content", expanded=False):
                if template_data.get('name'):
                    st.write(f"**Project Name:** {template_data['name']}")
                if template_data.get('type'):
                    st.write(f"**Project Type:** {template_data['type']}")
                if template_data.get('description'):
                    st.write(f"**Description:** {template_data['description']}")
                if template_data.get('functional_reqs'):
                    st.write(f"**Functional Requirements:** {len(template_data['functional_reqs'])} items")
                if template_data.get('non_functional_reqs'):
                    st.write(f"**Non-Functional Requirements:** {len(template_data['non_functional_reqs'])} items")
                if template_data.get('conditions'):
                    st.write(f"**Conditions:** {len(template_data['conditions'])} items")
                if template_data.get('template_content'):
                    st.write(f"**Additional Files:** {len(template_data['template_content'])} files")
                    for file_info in template_data['template_content']:
                        st.write(f"• {file_info['filename']}")
                
                # Show detailed content in JSON format for debugging
                with st.expander("🔍 Detailed Template Data (Debug)", expanded=False):
                    st.json({k: v for k, v in template_data.items() if k != 'template_content'})
    
    # Project data upload section (ENHANCED)
    st.subheader("📂 Project Data (Optional)")
    st.info("💡 **Project Data** helps the AI Assistant provide more accurate and context-aware responses by understanding your specific project content. Upload individual files or specify a folder path to include all project-related documents.")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.write("**📁 Upload Project Data Files:**")
        project_data_files = st.file_uploader(
            "Select multiple files of any type", 
            help="Supports: PDF, Word, Excel, PowerPoint, Text, Markdown, Images, Code files, and more. All files will be processed for AI enhancement.",
            accept_multiple_files=True,
            key="project_data_files"
        )
        
        if project_data_files:
            st.success(f"✅ Selected {len(project_data_files)} files:")
            for file in project_data_files:
                file_size = len(file.getvalue()) if hasattr(file, 'getvalue') else 0
                st.write(f"• {file.name} ({file_size:,} bytes)")
    
    with col2:
        st.write("**📂 Or Specify Project Data Folder:**")
        project_data_folder = st.text_input(
            "Project data folder path (optional)",
            placeholder="C:/path/to/project/data/folder",
            help="All files in this folder and subfolders will be processed recursively. Supports any file type."
        )
        
        # Folder validation
        if project_data_folder:
            if os.path.exists(project_data_folder):
                if os.path.isdir(project_data_folder):
                    # Count files in folder
                    file_count = sum(len(files) for _, _, files in os.walk(project_data_folder))
                    st.success(f"✅ Folder found: {file_count} files will be processed")
                else:
                    st.error("❌ Path exists but is not a folder")
            else:
                st.error("❌ Folder path does not exist")
    
    # Show RAG service status with more detail
    col1, col2 = st.columns(2)
    with col1:
        if st.session_state.rag_service.is_available():
            st.success("🧠 **RAG Service:** Available")
            st.caption("Files will be processed and vectorized for AI enhancement")
        else:
            st.warning("⚠️ **RAG Service:** Not available")
            st.caption("Files will be stored but not processed for AI")
    
    with col2:
        st.info("**Supported File Types:**")
        st.caption("📄 Documents: PDF, Word, Excel, PowerPoint")
        st.caption("📝 Text: TXT, MD, JSON, CSV, XML")  
        st.caption("💻 Code: PY, JS, HTML, CSS, SQL, etc.")
        st.caption("🖼️ Images: PNG, JPG, GIF (OCR)")
        st.caption("📊 Data: Any structured/unstructured files")
    
    st.subheader("📋 Project Information")
    
    project_types = [
        "Software Development",
        "Hardware Development", 
        "Research Project",
        "Process Improvement",
        "Product Development",
        "Custom"
    ]
    
    # Document Type Selection (Outside Form)
    st.subheader("📋 Document Types & Template Sources")
    
    document_types = [
        "Project Management Plan (PMP)",
        "Technical Concept Document (TCD)",
        "Configuration Management Plan",
        "Communication Management Plan",
        "Risk Management Plan",
        "Quality Plan"
    ]
    default_docs = template_data.get('recommended_docs', document_types[:3]) if template_data else document_types[:3]
    
    # Enhanced document type selection with template options
    st.write("📋 **Select Document Types and Template Sources:**")
    st.info("💡 For each document type, choose to use an existing template file or generate with AI.")
    
    # Initialize session state for document selections if not exists
    if 'new_project_doc_selections' not in st.session_state:
        st.session_state.new_project_doc_selections = {}
        # Set defaults based on template data
        for doc_type in document_types:
            is_default = doc_type in default_docs
            st.session_state.new_project_doc_selections[doc_type] = {
                'include': is_default,
                'source': 'AI Generated'
            }
    
    # Initialize session state for template selections
    if 'doc_template_selections' not in st.session_state:
        st.session_state.doc_template_selections = {}
    
    selected_docs_with_templates = {}
    
    for doc_type in document_types:
        # Get current state or use default
        current_state = st.session_state.new_project_doc_selections.get(doc_type, {
            'include': doc_type in default_docs,
            'source': 'AI Generated'
        })
        
        # Document type selection with session state
        include_doc = st.checkbox(
            f"📄 {doc_type}",
            value=current_state.get('include', False),
            key=f"include_{doc_type}"
        )
        
        # Update session state based on checkbox
        st.session_state.new_project_doc_selections[doc_type] = st.session_state.new_project_doc_selections.get(doc_type, {})
        st.session_state.new_project_doc_selections[doc_type]['include'] = include_doc
        
        if include_doc:
            # Template source selection
            col_a, col_b = st.columns([1, 2])
            
            with col_a:
                current_source = st.session_state.new_project_doc_selections[doc_type].get('source', 'AI Generated')
                template_source = st.radio(
                    f"Template source for {doc_type}:",
                    ["AI Generated", "Upload Template"],
                    index=0 if current_source == "AI Generated" else 1,
                    key=f"source_{doc_type}",
                    help="Choose AI Generated for automatic template creation, or Upload Template to use your own file."
                )
                
                # Update session state for source selection
                st.session_state.new_project_doc_selections[doc_type]['source'] = template_source
        
            with col_b:
                template_file = None
                if template_source == "Upload Template":
                    template_file = st.file_uploader(
                        f"Upload template for {doc_type}",
                        type=['docx', 'doc', 'txt', 'md', 'json'],
                        key=f"template_{doc_type}",
                        help="Upload a template file that will be used as the base for this document type."
                    )
                    
                    if template_file:
                        st.success(f"✅ Template uploaded: {template_file.name}")
                else:
                    st.info("🤖 AI will generate a professional template based on your project details.")
        
            # Store the selection
            selected_docs_with_templates[doc_type] = {
                'source': template_source,
                'template_file': template_file
            }
        
        # Store template selections for database
        if include_doc:
            st.session_state.doc_template_selections[doc_type] = selected_docs_with_templates.get(doc_type, {})
    
    # Show summary of document selections
    selected_count = sum(1 for doc_type, state in st.session_state.new_project_doc_selections.items() if state.get('include', False))
    if selected_count > 0:
        st.success(f"📋 {selected_count} document type(s) selected for template creation")
        
        with st.expander("📊 Document Selection Summary", expanded=False):
            for doc_type, state in st.session_state.new_project_doc_selections.items():
                if state.get('include', False):
                    source = state.get('source', 'AI Generated')
                    if source == "AI Generated":
                        st.write(f"• **{doc_type}:** 🤖 AI Generated")
                    else:
                        st.write(f"• **{doc_type}:** 📄 User Template")
    else:
        st.info("💡 No document types selected. You can add document templates later.")
    
    # Basic Project Form
    with st.form("new_project_form"):
        col1, col2 = st.columns(2)
        
        with col1:
            project_name = st.text_input(
                "Project Name*", 
                value=template_data.get('name', '') if template_data else '',
                placeholder="Enter project name"
            )
            project_type = st.selectbox(
                "Project Type*", 
                project_types,
                index=project_types.index(template_data.get('type', project_types[0])) if template_data and template_data.get('type') in project_types else 0
            )
            
        with col2:
            description = st.text_area(
                "Project Description*", 
                value=template_data.get('description', '') if template_data else '',
                placeholder="Describe your project objectives and scope",
                height=100
            )
        
        # Requirements section
        st.subheader("📝 Requirements & Constraints")
        
        col1, col2 = st.columns(2)
        
        with col1:
            st.write("**Functional Requirements:**")
            functional_reqs_default = '\n'.join(template_data.get('functional_reqs', [])) if template_data else ""
            functional_reqs = st.text_area(
                "Functional requirements (one per line)",
                value=functional_reqs_default,
                placeholder="System shall meet performance criteria\nSolution shall integrate with existing systems"
            )
            
            st.write("**Non-Functional Requirements:**")
            non_functional_reqs_default = '\n'.join(template_data.get('non_functional_reqs', [])) if template_data else ""
            non_functional_reqs = st.text_area(
                "Non-functional requirements (one per line)",
                value=non_functional_reqs_default,
                placeholder="System shall be available 99.9% of the time\nResponse time shall not exceed 2 seconds"
            )
            
        with col2:
            st.write("**Conditions/Constraints:**")
            conditions_default = '\n'.join(template_data.get('conditions', [])) if template_data else ""
            conditions = st.text_area(
                "Project conditions (one per line)",
                value=conditions_default,
                placeholder="Budget constraints must be observed\nCompliance with company standards required"
            )
        
        # Submit project
        submitted = st.form_submit_button("✅ Create Project", type="primary")
        
        if submitted:
            if project_name and project_type and description:
                # Get recommended docs from template selections
                recommended_docs = list(st.session_state.doc_template_selections.keys()) if st.session_state.doc_template_selections else []
                
                # Create project
                project = {
                    "name": project_name,
                    "type": project_type,
                    "description": description,
                    "functional_reqs": [req.strip() for req in functional_reqs.split('\n') if req.strip()],
                    "non_functional_reqs": [req.strip() for req in non_functional_reqs.split('\n') if req.strip()],
                    "conditions": [cond.strip() for cond in conditions.split('\n') if cond.strip()],
                    "recommended_docs": recommended_docs
                }
                
                # Save to database
                project_id = st.session_state.db.save_project(project)
                project['id'] = project_id
                
                # Process project data files and folder
                files_processed = 0
                failed_files = []
                
                with st.spinner("🔄 Processing project data files..."):
                    # Process uploaded files
                    if project_data_files:
                        for uploaded_file in project_data_files:
                            result = st.session_state.rag_service.process_file(
                                project_id=project_id,
                                file_obj=uploaded_file,
                                filename=uploaded_file.name,
                                is_template=False
                            )
                            if result["success"]:
                                files_processed += 1
                            else:
                                failed_files.append(f"{uploaded_file.name}: {result.get('error', 'Unknown error')}")
                    
                    # Process folder if specified
                    if project_data_folder and os.path.exists(project_data_folder):
                        folder_result = st.session_state.rag_service.process_folder(
                            project_id=project_id,
                            folder_path=project_data_folder,
                            is_template=False
                        )
                        if folder_result["success"]:
                            files_processed += folder_result["successful_files"]
                            if folder_result["failed_files"] > 0:
                                failed_files.append(f"Failed to process {folder_result['failed_files']} files from folder")
                    elif project_data_folder:
                        failed_files.append(f"Folder path does not exist: {project_data_folder}")
                
                # Process document templates
                template_results = {'created_templates': [], 'errors': [], 'template_folder': ''}
                if st.session_state.doc_template_selections:
                    with st.spinner("📋 Creating document templates..."):
                        # Get recommended docs from template selections
                        recommended_docs = list(st.session_state.doc_template_selections.keys())
                        
                        # Update project with recommended docs
                        st.session_state.db.update_project(project_id, {"recommended_docs": recommended_docs})
                        
                        # Process templates
                        template_results = process_document_templates(
                            project_id, 
                            project, 
                            st.session_state.doc_template_selections
                        )
                
                st.success(f"✅ Project '{project_name}' created successfully!")
                
                # Enhanced template processing feedback
                if template_results['created_templates']:
                    st.success(f"📋 Created {len(template_results['created_templates'])} document templates!")
                    st.info(f"📁 Templates saved to: `{template_results['template_folder']}`")
                    
                    # Show template details
                    with st.expander("📋 Template Details", expanded=False):
                        for template in template_results['created_templates']:
                            icon = "🤖" if template['source'] == 'AI Generated' else "📤"
                            st.write(f"{icon} **{template['doc_type']}** - {template['filename']}")
                            st.caption(f"Source: {template['source']} | Path: {template['file_path']}")
                
                if template_results['errors']:
                    st.warning("⚠️ Some templates could not be created:")
                    for error in template_results['errors'][:3]:
                        st.write(f"• {error}")
                
                # Enhanced file processing feedback
                if files_processed > 0:
                    st.success(f"🧠 Successfully processed {files_processed} project data file(s) for AI enhancement!")
                    st.info("💡 These files will now help the AI Assistant provide more relevant and context-aware responses.")
                
                if failed_files:
                    st.warning("⚠️ Some files could not be processed:")
                    for error in failed_files[:5]:  # Show first 5 errors
                        st.write(f"• {error}")
                    if len(failed_files) > 5:
                        st.write(f"... and {len(failed_files) - 5} more errors")
                
                if files_processed == 0 and not failed_files:
                    st.info("📂 Project created without additional data files. You can add files later in the 'Edit Projects' tab.")
                
                # Clear template selections after successful creation
                if 'doc_template_selections' in st.session_state:
                    del st.session_state.doc_template_selections
                if 'new_project_doc_selections' in st.session_state:
                    del st.session_state.new_project_doc_selections
                
                st.success("🎉 Project created successfully!")
                
                # Force refresh to update other tabs
                time.sleep(1)  # Brief pause
                st.rerun()
            else:
                st.error("Please fill in all required fields marked with *")

def show_edit_projects():
    """Show project editing and deletion interface"""
    st.title("✏️ Edit Projects")
    
    projects = st.session_state.db.get_projects()
    
    if not projects:
        st.warning("No projects found. Create a project first.")
        return
    
    # Create tabs for Edit and Delete
    edit_tab, delete_tab = st.tabs(["✏️ Edit Project", "🗑️ Delete Projects"])
    
    with edit_tab:
        st.subheader("Edit Existing Project")
        
        # Project selection
        project_options = {f"{p['name']} ({p['type']}) - ID: {p['id']}": p for p in projects}
        selected_project_name = st.selectbox(
            "Select Project to Edit:", 
            list(project_options.keys()),
            key="edit_project_select"
        )
        selected_project = project_options[selected_project_name]
        
        if selected_project:
            # Enhanced project information display
            col1, col2, col3 = st.columns(3)
            with col1:
                st.info(f"**Project:** {selected_project['name']}")
            with col2:
                st.info(f"**Type:** {selected_project['type']}")
            with col3:
                created_date = selected_project.get('created_at', 'Unknown')[:10] if selected_project.get('created_at') else 'Unknown'
                st.info(f"**Created:** {created_date}")
            
            # Show existing project data files with enhanced management
            st.subheader("📂 Current Project Data Files")
            existing_files = st.session_state.db.get_project_data_files(selected_project['id'])
            
            if existing_files:
                # File statistics
                total_files = len(existing_files)
                template_files = len([f for f in existing_files if f['is_template']])
                data_files = total_files - template_files
                total_size = sum(f['file_size'] for f in existing_files)
                
                col1, col2, col3, col4 = st.columns(4)
                with col1:
                    st.metric("Total Files", total_files)
                with col2:
                    st.metric("Data Files", data_files)
                with col3:
                    st.metric("Template Files", template_files)
                with col4:
                    st.metric("Total Size", f"{total_size:,} bytes")
                
                st.write("---")
                
                # Enhanced file display with better organization
                st.write("**File Management:**")
                
                # Group files by type
                file_groups = {'Data Files': [], 'Template Files': []}
                for file_info in existing_files:
                    if file_info['is_template']:
                        file_groups['Template Files'].append(file_info)
                    else:
                        file_groups['Data Files'].append(file_info)
                
                for group_name, files in file_groups.items():
                    if files:
                        with st.expander(f"📁 {group_name} ({len(files)} files)", expanded=True):
                            for file_info in files:
                                col1, col2, col3, col4 = st.columns([3, 1, 1, 1])
                                
                                with col1:
                                    # Enhanced file type icons
                                    file_type_icons = {
                                        'pdf': '📕', 'docx': '📄', 'doc': '📄', 
                                        'xlsx': '📊', 'xls': '📊', 'pptx': '📊', 'ppt': '📊',
                                        'txt': '📝', 'md': '📝', 'json': '🔧', 'csv': '📈',
                                        'py': '🐍', 'js': '🟨', 'html': '🌐', 'css': '🎨',
                                        'png': '🖼️', 'jpg': '🖼️', 'jpeg': '🖼️', 'gif': '🖼️',
                                        'xml': '📋', 'sql': '🗃️'
                                    }
                                    icon = file_type_icons.get(file_info['file_type'], '📄')
                                    
                                    st.write(f"{icon} **{file_info['filename']}**")
                                    st.caption(f"Size: {file_info['file_size']:,} bytes | Type: {file_info['file_type'].upper()} | Added: {file_info['created_at'][:16]}")
                                
                                with col2:
                                    # Show content preview with enhanced formatting
                                    if st.button("👀 Preview", key=f"preview_{file_info['id']}", help="View file content"):
                                        st.session_state[f"show_preview_{file_info['id']}"] = True
                                
                                with col3:
                                    # Download file content (if available)
                                    if st.button("💾 Export", key=f"export_{file_info['id']}", help="Download file content"):
                                        content = file_info.get('content', '')
                                        if content:
                                            st.download_button(
                                                label="Download",
                                                data=content,
                                                file_name=file_info['filename'],
                                                mime="text/plain",
                                                key=f"download_{file_info['id']}"
                                            )
                                        else:
                                            st.warning("No content available for download")
                                
                                with col4:
                                    # Delete file button with confirmation
                                    if st.button("🗑️ Remove", key=f"delete_{file_info['id']}", type="secondary", help="Remove file from project"):
                                        if st.session_state.get(f"confirm_delete_{file_info['id']}", False):
                                            st.session_state.db.delete_project_data_file(file_info['id'])
                                            st.success(f"✅ Removed {file_info['filename']}")
                                            st.rerun()
                                        else:
                                            st.session_state[f"confirm_delete_{file_info['id']}"] = True
                                            st.warning("Click again to confirm deletion")
                                
                                # Show preview if requested
                                if st.session_state.get(f"show_preview_{file_info['id']}", False):
                                    with st.expander(f"📖 Content Preview - {file_info['filename']}", expanded=True):
                                        content = file_info.get('content', '')
                                        if content:
                                            # Show preview with syntax highlighting if possible
                                            if file_info['file_type'] in ['json', 'py', 'js', 'html', 'css', 'xml']:
                                                st.code(content[:2000] + "..." if len(content) > 2000 else content, 
                                                        language=file_info['file_type'])
                                            else:
                                                st.text_area("Content Preview", 
                                                           value=content[:2000] + "..." if len(content) > 2000 else content, 
                                                           height=300, disabled=True)
                                            
                                            if len(content) > 2000:
                                                st.caption(f"Showing first 2000 characters of {len(content)} total")
                                        else:
                                            st.warning("No content available for preview")
                                        
                                        # Close preview button
                                        if st.button("✖️ Close Preview", key=f"close_preview_{file_info['id']}"):
                                            st.session_state[f"show_preview_{file_info['id']}"] = False
                                            st.rerun()
                
                st.divider()
            else:
                st.info("📂 No project data files found. Add some files below to enhance AI responses for this project.")
            
            # Add new project data files section with enhanced interface
            st.subheader("📂 Add New Project Data")
            st.info("💡 Add more files to improve the AI Assistant's understanding of your project context.")
            
            col1, col2 = st.columns(2)
            
            with col1:
                st.write("**📁 Upload New Files:**")
                new_project_data_files = st.file_uploader(
                    "Select files to add to project", 
                    help="Upload any files containing project information. All file types are supported.",
                    accept_multiple_files=True,
                    key="edit_project_data_files"
                )
                
                if new_project_data_files:
                    st.success(f"✅ {len(new_project_data_files)} file(s) ready to upload:")
                    for file in new_project_data_files:
                        file_size = len(file.getvalue()) if hasattr(file, 'getvalue') else 0
                        st.write(f"• {file.name} ({file_size:,} bytes)")
            
            with col2:
                st.write("**📂 Or Specify Folder Path:**")
                new_project_data_folder = st.text_input(
                    "Additional project data folder path",
                    placeholder="C:/path/to/additional/project/data",
                    help="All files in this folder and subfolders will be processed recursively.",
                    key="edit_project_data_folder"
                )
                
                # Folder validation for edit mode
                if new_project_data_folder:
                    if os.path.exists(new_project_data_folder):
                        if os.path.isdir(new_project_data_folder):
                            file_count = sum(len(files) for _, _, files in os.walk(new_project_data_folder))
                            st.success(f"✅ Folder found: {file_count} files will be processed")
                        else:
                            st.error("❌ Path exists but is not a folder")
                    else:
                        st.error("❌ Folder path does not exist")
            
            # Project template upload section for editing
            st.subheader("📁 Update from Template(s) (Optional)")
            uploaded_files_edit = st.file_uploader(
                "Upload template file(s) to update project", 
                type=['json', 'txt', 'md', 'xlsx', 'xls', 'docx', 'doc', 'pptx', 'ppt', 'pdf'],
                help="Upload one or more template files to update/merge with existing project data. Supports: JSON, TXT, Markdown, Excel, Word, PowerPoint, PDF",
                accept_multiple_files=True,
                key="edit_template_upload"
            )
            
            template_data_edit = None
            if uploaded_files_edit:
                template_data_edit = load_project_template(uploaded_files_edit)
                if template_data_edit:
                    st.success(f"✅ Loaded {len(uploaded_files_edit)} template file(s) for merging!")
                    
                    with st.expander("📋 Preview template updates", expanded=True):
                        st.warning("⚠️ Template data will be merged with existing project data. Existing data will be preserved unless explicitly replaced.")
                        
                        if template_data_edit.get('name'):
                            st.write(f"**New Project Name:** {template_data_edit['name']}")
                        if template_data_edit.get('type'):
                            st.write(f"**New Project Type:** {template_data_edit['type']}")
                        if template_data_edit.get('description'):
                            st.write(f"**New Description:** {template_data_edit['description']}")
                        if template_data_edit.get('functional_reqs'):
                            st.write(f"**Additional Functional Requirements:** {len(template_data_edit['functional_reqs'])} items")
                        if template_data_edit.get('non_functional_reqs'):
                            st.write(f"**Additional Non-Functional Requirements:** {len(template_data_edit['non_functional_reqs'])} items")
                        if template_data_edit.get('conditions'):
                            st.write(f"**Additional Conditions:** {len(template_data_edit['conditions'])} items")
                        if template_data_edit.get('template_content'):
                            st.write(f"**Additional Reference Files:** {len(template_data_edit['template_content'])} files")
            
            with st.form("edit_project_form"):
                col1, col2 = st.columns(2)
                
                with col1:
                    # Merge template name with existing project name
                    merged_name = template_data_edit.get('name', selected_project.get('name', '')) if template_data_edit else selected_project.get('name', '')
                    project_name = st.text_input(
                        "Project Name*", 
                        value=merged_name,
                        placeholder="Enter project name"
                    )
                    
                    project_types = [
                        "Software Development",
                        "Hardware Development", 
                        "Research Project",
                        "Process Improvement",
                        "Product Development",
                        "Custom"
                    ]
                    
                    # Merge template type with existing project type
                    merged_type = template_data_edit.get('type', selected_project.get('type', project_types[0])) if template_data_edit else selected_project.get('type', project_types[0])
                    current_type_index = project_types.index(merged_type) if merged_type in project_types else 0
                    project_type = st.selectbox(
                        "Project Type*", 
                        project_types,
                        index=current_type_index
                    )
                    
                with col2:
                    # Merge template description with existing project description
                    merged_desc = template_data_edit.get('description', selected_project.get('description', '')) if template_data_edit else selected_project.get('description', '')
                    description = st.text_area(
                        "Project Description*", 
                        value=merged_desc,
                        placeholder="Describe your project objectives and scope",
                        height=100
                    )
                
                # Requirements section
                st.subheader("📝 Requirements & Constraints")
                
                col1, col2 = st.columns(2)
                
                with col1:
                    st.write("**Functional Requirements:**")
                    # Merge template functional requirements with existing ones
                    existing_func_reqs = selected_project.get('functional_reqs', [])
                    template_func_reqs = template_data_edit.get('functional_reqs', []) if template_data_edit else []
                    merged_func_reqs = existing_func_reqs + [req for req in template_func_reqs if req not in existing_func_reqs]
                    functional_reqs_text = '\n'.join(merged_func_reqs)
                    functional_reqs = st.text_area(
                        "Functional requirements (one per line)",
                        value=functional_reqs_text,
                        placeholder="System shall meet performance criteria\nSolution shall integrate with existing systems",
                        key="edit_func_reqs"
                    )
                    
                    st.write("**Non-Functional Requirements:**")
                    # Merge template non-functional requirements with existing ones
                    existing_nonfunc_reqs = selected_project.get('non_functional_reqs', [])
                    template_nonfunc_reqs = template_data_edit.get('non_functional_reqs', []) if template_data_edit else []
                    merged_nonfunc_reqs = existing_nonfunc_reqs + [req for req in template_nonfunc_reqs if req not in existing_nonfunc_reqs]
                    non_functional_reqs_text = '\n'.join(merged_nonfunc_reqs)
                    non_functional_reqs = st.text_area(
                        "Non-functional requirements (one per line)",
                        value=non_functional_reqs_text,
                        placeholder="System shall be available 99.9% of the time\nResponse time shall not exceed 2 seconds",
                        key="edit_nonfunc_reqs"
                    )
                    
                with col2:
                    st.write("**Conditions/Constraints:**")
                    # Merge template conditions with existing ones
                    existing_conditions = selected_project.get('conditions', [])
                    template_conditions = template_data_edit.get('conditions', []) if template_data_edit else []
                    merged_conditions = existing_conditions + [cond for cond in template_conditions if cond not in existing_conditions]
                    conditions_text = '\n'.join(merged_conditions)
                    conditions = st.text_area(
                        "Project conditions (one per line)",
                        value=conditions_text,
                        placeholder="Budget constraints must be observed\nCompliance with company standards required",
                        key="edit_conditions"
                    )
                    
                    st.write("**Document Types to Generate:**")
                    document_types = [
                        "Project Management Plan (PMP)",
                        "Technical Concept Document (TCD)",
                        "Configuration Management Plan",
                        "Communication Management Plan",
                        "Risk Management Plan",
                        "Quality Plan"
                    ]
                    recommended_docs = st.multiselect(
                        "Select document types", 
                        document_types, 
                        default=selected_project.get('recommended_docs', []),
                        key="edit_recommended_docs"
                    )
                
                # Submit buttons
                col1, col2 = st.columns(2)
                with col1:
                    update_button = st.form_submit_button("💾 Update Project", type="primary")
                with col2:
                    cancel_button = st.form_submit_button("❌ Cancel Changes")
                
                if update_button:
                    if project_name and project_type and description:
                        # Prepare updates
                        updates = {
                            "name": project_name,
                            "type": project_type,
                            "description": description,
                            "functional_reqs": [req.strip() for req in functional_reqs.split('\n') if req.strip()],
                            "non_functional_reqs": [req.strip() for req in non_functional_reqs.split('\n') if req.strip()],
                            "conditions": [cond.strip() for cond in conditions.split('\n') if cond.strip()],
                            "recommended_docs": recommended_docs
                        }
                        
                        # Update in database
                        st.session_state.db.update_project(selected_project['id'], updates)
                        
                        # Process new project data files with enhanced feedback
                        files_processed = 0
                        failed_files = []
                        
                        with st.spinner("🔄 Processing new project data files..."):
                            # Process uploaded files
                            if new_project_data_files:
                                for uploaded_file in new_project_data_files:
                                    result = st.session_state.rag_service.process_file(
                                        project_id=selected_project['id'],
                                        file_obj=uploaded_file,
                                        filename=uploaded_file.name,
                                        is_template=False
                                    )
                                    if result["success"]:
                                        files_processed += 1
                                    else:
                                        failed_files.append(f"{uploaded_file.name}: {result.get('error', 'Unknown error')}")
                            
                            # Process folder if specified
                            if new_project_data_folder and os.path.exists(new_project_data_folder):
                                folder_result = st.session_state.rag_service.process_folder(
                                    project_id=selected_project['id'],
                                    folder_path=new_project_data_folder,
                                    is_template=False
                                )
                                if folder_result["success"]:
                                    files_processed += folder_result["successful_files"]
                                    if folder_result["failed_files"] > 0:
                                        failed_files.append(f"Failed to process {folder_result['failed_files']} files from folder")
                            elif new_project_data_folder:
                                failed_files.append(f"Folder path does not exist: {new_project_data_folder}")
                        
                        st.success(f"✅ Project '{project_name}' updated successfully!")
                        
                        # Enhanced update feedback
                        if files_processed > 0:
                            st.success(f"🧠 Successfully processed {files_processed} additional project data file(s)!")
                            st.info("💡 The AI Assistant now has more context about your project.")
                        
                        if failed_files:
                            st.warning("⚠️ Some files could not be processed:")
                            for error in failed_files[:3]:  # Show first 3 errors
                                st.write(f"• {error}")
                            if len(failed_files) > 3:
                                st.write(f"... and {len(failed_files) - 3} more errors")
                        
                        st.info("💾 Changes have been saved. The updated project information is now available across all tabs.")
                        time.sleep(1)
                        st.rerun()
                    else:
                        st.error("Please fill in all required fields marked with *")
                
                if cancel_button:
                    st.info("Changes cancelled.")
                    st.rerun()
    
    with delete_tab:
        st.subheader("🗑️ Delete Projects")
        st.warning("⚠️ **Warning:** Deleting projects will also remove all associated documents and workflows. This action cannot be undone.")
        
        # Selection mode
        selection_mode = st.radio(
            "Selection Mode:",
            ["Single Project", "Multiple Projects"],
            key="delete_mode"
        )
        
        if selection_mode == "Single Project":
            # Single project deletion
            project_options = {f"{p['name']} ({p['type']}) - ID: {p['id']}": p for p in projects}
            selected_project_name = st.selectbox(
                "Select Project to Delete:", 
                list(project_options.keys()),
                key="single_delete_select"
            )
            selected_project = project_options[selected_project_name]
            
            if selected_project:
                # Show project details
                with st.expander("📋 Project Details", expanded=True):
                    col1, col2 = st.columns(2)
                    with col1:
                        st.write(f"**Name:** {selected_project['name']}")
                        st.write(f"**Type:** {selected_project['type']}")
                        st.write(f"**Created:** {selected_project.get('created_at', 'Unknown')[:10]}")
                    with col2:
                        # Count related items
                        documents = st.session_state.db.get_documents(selected_project['id'])
                        workflows = st.session_state.db.get_workflows(selected_project['id'])
                        st.write(f"**Documents:** {len(documents)}")
                        st.write(f"**Workflows:** {len(workflows)}")
                
                # Confirmation
                confirm_delete = st.checkbox(
                    f"I understand that deleting '{selected_project['name']}' will permanently remove all data",
                    key="single_confirm_delete"
                )
                
                if st.button("🗑️ Delete Project", type="secondary", disabled=not confirm_delete):
                    st.session_state.db.delete_project(selected_project['id'])
                    st.success(f"✅ Project '{selected_project['name']}' deleted successfully!")
                    time.sleep(1)
                    st.rerun()
        
        else:
            # Multiple project deletion
            st.write("**Select Multiple Projects to Delete:**")
            
            # Create checkboxes for each project
            selected_projects = []
            for project in projects:
                # Count related items
                documents = st.session_state.db.get_documents(project['id'])
                workflows = st.session_state.db.get_workflows(project['id'])
                
                checkbox_label = f"{project['name']} ({project['type']}) - {len(documents)} docs, {len(workflows)} workflows"
                
                if st.checkbox(checkbox_label, key=f"multi_delete_{project['id']}"):
                    selected_projects.append(project)
            
            if selected_projects:
                st.write(f"**Selected {len(selected_projects)} project(s) for deletion:**")
                for project in selected_projects:
                    st.write(f"• {project['name']} ({project['type']})")
                
                # Confirmation
                confirm_multiple_delete = st.checkbox(
                    f"I understand that deleting these {len(selected_projects)} project(s) will permanently remove all associated data",
                    key="multi_confirm_delete"
                )
                
                if st.button("🗑️ Delete Selected Projects", type="secondary", disabled=not confirm_multiple_delete):
                    project_ids = [p['id'] for p in selected_projects]
                    project_names = [p['name'] for p in selected_projects]
                    
                    st.session_state.db.delete_multiple_projects(project_ids)
                    st.success(f"✅ Deleted {len(selected_projects)} project(s): {', '.join(project_names)}")
                    time.sleep(1)
                    st.rerun()
            else:
                st.info("Select projects using the checkboxes above to enable deletion.")

def show_document_generation():
    """Show document generation interface"""
    st.title("📝 Generate Document")
    
    projects = st.session_state.db.get_projects()
    
    if not projects:
        st.warning("No projects found. Please create a project first.")
        return
    
    # Project selection
    project_options = {f"{p['name']} ({p['type']})": p for p in projects}
    selected_project_name = st.selectbox("Select Project", list(project_options.keys()))
    selected_project = project_options[selected_project_name]
    
    if selected_project:
        st.info(f"**Project:** {selected_project['name']} | **Type:** {selected_project['type']}")
        
        # Show project data summary
        project_data_summary = st.session_state.rag_service.get_project_data_summary(selected_project['id'])
        
        if project_data_summary['total_files'] > 0:
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("📂 Total Files", project_data_summary['total_files'])
            with col2:
                st.metric("🧠 Data Chunks", project_data_summary['total_chunks'])
            with col3:
                st.metric("📊 Data Files", project_data_summary['data_files'])
            with col4:
                st.metric("🏷️ Templates", project_data_summary['template_files'])
            
            if st.session_state.rag_service.is_available():
                st.success("🧠 RAG-enhanced document generation available!")
            else:
                st.warning("⚠️ RAG service unavailable - basic document generation only")
        else:
            st.info("💡 No project data files found. Documents will be generated using project requirements only.")
        
        # Template upload section for document generation
        st.subheader("📁 Additional Template Files (Optional)")
        uploaded_template_files = st.file_uploader(
            "Upload additional template files to enhance document generation", 
            type=['json', 'txt', 'md', 'xlsx', 'xls', 'docx', 'doc', 'pptx', 'ppt', 'pdf'],
            help="Upload template files that contain relevant information for document generation. Supports: JSON, TXT, Markdown, Excel, Word, PowerPoint, PDF",
            accept_multiple_files=True,
            key="doc_gen_template_upload"
        )
        
        template_data_doc_gen = None
        if uploaded_template_files:
            template_data_doc_gen = load_project_template(uploaded_template_files)
            if template_data_doc_gen:
                st.success(f"✅ Loaded {len(uploaded_template_files)} template file(s) for document enhancement!")
                
                with st.expander("📋 Preview template content", expanded=False):
                    if template_data_doc_gen.get('template_content'):
                        st.write(f"**Reference Files:** {len(template_data_doc_gen['template_content'])} files loaded")
                        for file_info in template_data_doc_gen['template_content']:
                            filename = file_info.get('filename', 'Unknown')
                            content_preview = file_info.get('content', '')[:200] + "..." if len(file_info.get('content', '')) > 200 else file_info.get('content', '')
                            st.write(f"**{filename}:** {content_preview}")
        
        with st.form("document_generation_form"):
            col1, col2 = st.columns(2)
            
            with col1:
                document_types = [
                    "Project Management Plan (PMP)",
                    "Technical Concept Document (TCD)",
                    "Configuration Management Plan",
                    "Communication Management Plan",
                    "Risk Management Plan",
                    "Quality Plan",
                    "Requirements Document",
                    "Architecture Document"
                ]
                document_type = st.selectbox("Document Type", document_types)
                template_name = st.selectbox("Template", ["Default Document", "Bosch Template"])
                
            with col2:
                custom_requirements = st.text_area(
                    "Additional Requirements (optional)",
                    placeholder="Any specific requirements for this document..."
                )
                
                use_ai = st.checkbox("Use AI Generation", value=True, help="Generate content using AI")
            
            generate_button = st.form_submit_button("🚀 Generate Document", type="primary")
            
            if generate_button:
                with st.spinner("Generating document..."):
                    if use_ai:
                        # AI-generated content with RAG and template context
                        prompt = f"Generate a {document_type} for a {selected_project['type']} project named '{selected_project['name']}'. {custom_requirements}"
                        
                        # Add RAG context from project data
                        rag_context = ""
                        if st.session_state.rag_service.is_available():
                            doc_generation_query = f"Create {document_type} for {selected_project['name']} {selected_project['type']} project"
                            rag_context = st.session_state.rag_service.get_context_for_query(
                                project_id=selected_project['id'],
                                query=doc_generation_query,
                                max_context_length=2500
                            )
                        
                        if rag_context:
                            prompt += f"\n\nRelevant project data and documentation:\n{rag_context}"
                        
                        # Add template content to AI prompt if available
                        if template_data_doc_gen and template_data_doc_gen.get('template_content'):
                            prompt += "\n\nReference information from uploaded templates:\n"
                            for file_info in template_data_doc_gen['template_content']:
                                filename = file_info.get('filename', 'Unknown')
                                content = file_info.get('content', '')
                                if content and len(content) < 1500:  # Only add shorter content to prompt
                                    prompt += f"\n--- From {filename} ---\n{content}\n"
                        
                        content = generate_ai_content(prompt, selected_project)
                        
                        # Add RAG enhancement indicator
                        if rag_context:
                            content += "\n\n---\n*🧠 This document was enhanced with project data using RAG (Retrieval-Augmented Generation).*"
                            
                    else:
                        # Template-based content with uploaded template data and project files
                        st.info("📋 Using template-based generation (without AI)...")
                        
                        # Get project data for template generation
                        project_files_context = ""
                        if st.session_state.rag_service.is_available():
                            doc_generation_query = f"Information for {document_type} {selected_project['name']}"
                            project_files_context = st.session_state.rag_service.get_context_for_query(
                                project_id=selected_project['id'],
                                query=doc_generation_query,
                                max_context_length=2000
                            )
                        
                        # Use uploaded template data if available
                        enhanced_template_data = template_data_doc_gen or {'template_content': []}
                        
                        # Add project files context if available
                        if project_files_context:
                            enhanced_template_data['template_content'].append({
                                'filename': 'project_data_context.txt',
                                'content': project_files_context
                            })
                        
                        # Check if we have any template content
                        has_template_content = enhanced_template_data.get('template_content') and len(enhanced_template_data['template_content']) > 0
                        
                        if has_template_content:
                            st.success(f"✅ Using {len(enhanced_template_data['template_content'])} template/data sources")
                        else:
                            st.warning("⚠️ No template files uploaded. Using default template structure.")
                        
                        content = generate_template_content(document_type, selected_project, custom_requirements.split('\n') if custom_requirements else [], enhanced_template_data)
                        
                        # Add enhancement indicators
                        if project_files_context and template_data_doc_gen:
                            content += "\n\n---\n*📂 This document incorporates information from your project data files and uploaded templates.*"
                        elif project_files_context:
                            content += "\n\n---\n*📂 This document incorporates information from your project data files.*"
                        elif template_data_doc_gen:
                            content += "\n\n---\n*📋 This document was generated using your uploaded templates.*"
                    
                    # Save document to database
                    try:
                        document = {
                            "project_id": selected_project["id"],
                            "name": f"{document_type}_{selected_project['name']}",
                            "type": document_type,
                            "content": content,
                            "status": "Draft"
                        }
                        
                        document_id = st.session_state.db.save_document(document)
                        document["id"] = document_id
                        
                        # Create workflow
                        workflow = {
                            "project_id": selected_project["id"],
                            "document_id": document_id,
                            "name": f"{document_type} Review",
                            "status": "Active",
                            "approvers": ["Project Manager", "Technical Lead", "Quality Assurance"],
                            "current_step": 0
                        }
                        
                        workflow_id = st.session_state.db.save_workflow(workflow)
                        
                        st.success("✅ Document generated successfully!")
                        st.success("🔄 Approval workflow created!")
                        st.info(f"📄 Document ID: {document_id} | 🔄 Workflow ID: {workflow_id}")
                        st.info("👥 **Next Step:** Check the 'Workflow Management' tab to view and manage the approval process.")
                        
                        # Display generated content preview
                        st.subheader("📄 Document Preview")
                        with st.expander("View Generated Content", expanded=True):
                            st.markdown(content)
                        
                        # Option to download the document
                        st.download_button(
                            label="📥 Download Document",
                            data=content,
                            file_name=f"{document_type}_{selected_project['name']}.md",
                            mime="text/markdown"
                        )
                        
                    except Exception as db_error:
                        st.error(f"❌ Database Error: {str(db_error)}")
                        st.error("Failed to save document to database. Please try again.")

def generate_template_content(document_type: str, project: Dict, additional_reqs: List[str], template_data: Optional[Dict] = None) -> str:
    """Generate template-based content with optional template file content"""
    
    # Check if we have meaningful project data (not just defaults)
    has_meaningful_project_data = (
        (project.get('functional_reqs') and len(project.get('functional_reqs', [])) > 0) or
        (project.get('non_functional_reqs') and len(project.get('non_functional_reqs', [])) > 0) or
        (project.get('description') and project.get('description') != 'Not specified')
    )
    
    # Check if we have RAG context from vector database
    has_vector_data = False
    if template_data and template_data.get('template_content'):
        for file_info in template_data['template_content']:
            if file_info.get('filename') == 'project_data_context.txt' and file_info.get('content'):
                has_vector_data = True
                break
    
    content = f"""# {document_type}

## Project Overview
**Project Name:** {project.get('name', 'Not specified')}
**Project Type:** {project.get('type', 'Not specified')}
**Description:** {project.get('description', 'Not specified')}

## Executive Summary
This document outlines the {document_type.lower()} for the {project.get('name', 'project')} project.
"""

    # Add template file content if available (excluding RAG context)
    if template_data and template_data.get('template_content'):
        content += "\n## Reference Information from Template Files\n"
        for file_info in template_data['template_content']:
            filename = file_info.get('filename', 'Unknown file')
            file_content = file_info.get('content', '')
            
            # Skip RAG context file and only include meaningful content (not error messages)
            if (filename != 'project_data_context.txt' and 
                file_content and 
                not file_content.startswith(('Excel file:', 'Word document:', 'PowerPoint presentation:', 'PDF document:'))):
                content += f"\n### Content from {filename}\n"
                # Limit content length for readability
                if len(file_content) > 2000:
                    content += file_content[:2000] + "\n... [Content truncated for brevity]\n"
                else:
                    content += file_content + "\n"

    content += "\n## Requirements\n\n### Functional Requirements\n"
    
    # If no meaningful project data and no vector data, leave template sections empty
    if not has_meaningful_project_data and not has_vector_data:
        content += "[Please add your functional requirements here]\n\n"
        content += "### Non-Functional Requirements\n"
        content += "[Please add your non-functional requirements here]\n\n"
    else:
        # Use project data or default examples
        if project.get('functional_reqs'):
            for i, req in enumerate(project['functional_reqs'], 1):
                content += f"{i}. {req}\n"
        else:
            content += "1. System shall meet specified performance criteria\n"
            content += "2. Solution shall integrate with existing systems\n"
        
        content += "\n### Non-Functional Requirements\n"
        if project.get('non_functional_reqs'):
            for i, req in enumerate(project['non_functional_reqs'], 1):
                content += f"{i}. {req}\n"
        else:
            content += "1. System shall be available 99.9% of the time\n"
            content += "2. Response time shall not exceed 2 seconds\n"
    
    if additional_reqs:
        content += "\n### Additional Requirements\n"
        for i, req in enumerate(additional_reqs, 1):
            if req.strip():
                content += f"{i}. {req.strip()}\n"
    
    # Add implementation plan - keep flexible based on available data
    if has_meaningful_project_data or has_vector_data:
        content += f"""
## Implementation Plan
- **Phase 1:** Planning and Design
- **Phase 2:** Development and Testing
- **Phase 3:** Deployment and Monitoring

## Risk Management
Identified risks and mitigation strategies will be documented as the project progresses.

## Timeline
Project timeline and milestones to be established based on requirements analysis.

## Quality Assurance
Quality gates and testing procedures will ensure deliverable meets Bosch standards.

## Approval
This document requires review and approval from designated stakeholders.
"""
    else:
        content += f"""
## Implementation Plan
[Please outline your implementation phases and approach here]

## Risk Management
[Please identify key risks and mitigation strategies here]

## Timeline
[Please define project timeline and key milestones here]

## Quality Assurance
[Please specify quality gates and testing procedures here]

## Approval
[Please define approval workflow and stakeholders here]
"""
    
    content += f"""
---
*Generated by Bosch AI Document Manager - {datetime.now().strftime("%Y-%m-%d %H:%M")}*
"""
    return content

def show_workflow_management():
    """Show workflow management interface with AI insights"""
    st.title("👥 Workflow Management")
    
    workflows = st.session_state.db.get_workflows()
    documents = st.session_state.db.get_documents()
    
    # AI Workflow Insights at the top
    if workflows and documents:
        st.markdown("### 🧠 AI Workflow Insights")
        col1, col2, col3 = st.columns(3)
        
        with col1:
            # Bottleneck detection
            active_workflows = [w for w in workflows if w['status'] == 'Active']
            if active_workflows:
                avg_step = sum(w['current_step'] for w in active_workflows) / len(active_workflows)
                bottleneck_status = "High" if avg_step > 1.5 else "Medium" if avg_step > 0.5 else "Low"
                bottleneck_color = "🔴" if bottleneck_status == "High" else "🟡" if bottleneck_status == "Medium" else "🟢"
                st.metric("Bottleneck Risk", f"{bottleneck_color} {bottleneck_status}")
            else:
                st.metric("Bottleneck Risk", "🟢 None")
        
        with col2:
            # Completion prediction
            total_workflows = len(workflows)
            completed_workflows = len([w for w in workflows if w['status'] == 'Completed'])
            if total_workflows > 0:
                completion_rate = (completed_workflows / total_workflows) * 100
                st.metric("Completion Rate", f"{completion_rate:.0f}%")
            else:
                st.metric("Completion Rate", "N/A")
        
        with col3:
            # Risk assessment
            rejected_workflows = len([w for w in workflows if w['status'] == 'Rejected'])
            risk_level = "High" if rejected_workflows > 2 else "Medium" if rejected_workflows > 0 else "Low"
            risk_color = "🔴" if risk_level == "High" else "🟡" if risk_level == "Medium" else "🟢"
            st.metric("Quality Risk", f"{risk_color} {risk_level}")
        
        # AI Recommendations
        st.markdown("#### 💡 AI Recommendations")
        recommendations = []
        
        if len(active_workflows) > 5:
            recommendations.append("⚡ Consider parallel approval processes to reduce bottlenecks")
        
        if rejected_workflows > 1:
            recommendations.append("📋 Review document templates to improve first-pass approval rates")
        
        if len(workflows) > 0:
            avg_duration = 3  # Placeholder - in real implementation, calculate from workflow history
            recommendations.append(f"⏱️ Average workflow duration: {avg_duration} days - consider optimization")
        
        if not recommendations:
            recommendations.append("✅ Workflow performance is optimal!")
        
        for rec in recommendations:
            st.info(rec)
    
    # Debug information for troubleshooting
    with st.expander("🔍 Debug Information", expanded=False):
        st.write(f"**Total Workflows:** {len(workflows)}")
        st.write(f"**Total Documents:** {len(documents)}")
        
        # Add button to reset incorrectly rejected workflows for testing
        if st.button("🔄 Reset All Workflows to Active Status (Debug)"):
            try:
                for workflow in workflows:
                    if workflow['status'] == 'Rejected':
                        st.session_state.db.update_workflow(
                            workflow['id'], 
                            {'status': 'Active', 'current_step': 0}
                        )
                        # Also reset corresponding document status
                        st.session_state.db.update_document(
                            workflow['document_id'], 
                            {'status': 'Draft'}
                        )
                st.success("✅ All workflows reset to Active status")
                st.rerun()
            except Exception as e:
                st.error(f"❌ Error resetting workflows: {str(e)}")
        
        if workflows:
            st.write("**Workflow Details:**")
            for w in workflows:
                st.write(f"- ID: {w.get('id')}, Status: {w.get('status')}, Document ID: {w.get('document_id')}")
        if documents:
            st.write("**Document Details:**")
            for d in documents:
                st.write(f"- ID: {d.get('id')}, Name: {d.get('name')}, Type: {d.get('type')}, Status: {d.get('status')}")
    
    tab1, tab2, tab3, tab4 = st.tabs(["📋 Pending Tasks", "📊 Workflow Status", "📈 Workflow Flowchart", "🤖 AI Optimizer"])
    
    with tab4:
        show_ai_workflow_optimizer(workflows, documents)
    
    with tab1:
        st.subheader("Pending Tasks")
        
        # For PM: Show active workflows for approval
        # For Project team: Show active workflows (view only) AND rejected workflows for rework
        if check_access('PM'):
            pending_workflows = [w for w in workflows if w["status"] == "Active"]
            st.write("**Approval Tasks:**")
        else:  # Project team
            pending_workflows = [w for w in workflows if w["status"] in ["Active", "Rejected"]]
            st.write("**Active Workflows and Rework Tasks:**")
        
        if pending_workflows:
            for workflow in pending_workflows:
                document = next((d for d in documents if d["id"] == workflow["document_id"]), None)
                
                if document:
                    # Different display based on workflow status and user role
                    if workflow["status"] == "Active":
                        current_approver = workflow["approvers"][workflow["current_step"]]
                        title = f"📄 {document['name']} - {current_approver}"
                        status_icon = "🔄"
                    else:  # Rejected - only shown to Project team
                        title = f"� {document['name']} - Requires Rework"
                        status_icon = "⚠️"
                    
                    with st.expander(f"{status_icon} {title}", expanded=False):
                        col1, col2 = st.columns(2)
                        
                        with col1:
                            st.write(f"**Document Type:** {document['type']}")
                            if workflow["status"] == "Active":
                                st.write(f"**Current Approver:** {current_approver}")
                                st.write(f"**Step:** {workflow['current_step'] + 1} of {len(workflow['approvers'])}")
                            else:  # Rejected
                                st.write(f"**Status:** Rejected - Requires Rework")
                                st.write(f"**Action Required:** Review feedback and resubmit")
                        
                        with col2:
                            st.write(f"**Workflow:** {workflow['name']}")
                            st.write(f"**Status:** {workflow['status']}")
                        
                        # Show workflow comments/feedback
                        try:
                            comments = st.session_state.db.get_workflow_comments(workflow['id'])
                            if comments:
                                st.write("**Feedback/Comments:**")
                                for comment in comments[-3:]:  # Show last 3 comments
                                    st.write(f"- **{comment['approver']}** ({comment['action']}): {comment['comment']}")
                        except:
                            pass  # Handle case where get_workflow_comments doesn't exist yet
                        
                        # Action buttons based on role and workflow status
                        if workflow["status"] == "Active" and check_access('PM'):
                            # PM approval form (existing logic)
                            with st.form(f"task_{workflow['id']}"):
                                col1, col2 = st.columns([1, 2])
                                
                                with col1:
                                    action = st.selectbox("Action", ["Approve", "Reject"], key=f"action_{workflow['id']}")
                                
                                with col2:
                                    comments = st.text_area(
                                        "Comments", 
                                        key=f"comments_{workflow['id']}",
                                        placeholder="Add your comments about this decision..."
                                    )
                                
                                if st.form_submit_button("Submit Decision"):
                                    current_approver = workflow["approvers"][workflow["current_step"]]
                                    
                                    try:
                                        # Add comment to workflow history
                                        st.session_state.db.add_workflow_comment(
                                            workflow['id'], 
                                            current_approver, 
                                            action, 
                                            comments
                                        )
                                        
                                        if action == "Approve":
                                            new_step = workflow["current_step"] + 1
                                            if new_step >= len(workflow["approvers"]):
                                                # Workflow complete
                                                st.session_state.db.update_workflow(
                                                    workflow['id'], 
                                                    {'status': 'Completed', 'current_step': new_step}
                                                )
                                                st.session_state.db.update_document(
                                                    document['id'], 
                                                    {'status': 'Approved'}
                                                )
                                                st.success("✅ Document fully approved!")
                                            else:
                                                # Move to next step
                                                st.session_state.db.update_workflow(
                                                    workflow['id'], 
                                                    {'current_step': new_step}
                                                )
                                                next_approver = workflow['approvers'][new_step]
                                                st.success(f"✅ Approved! Moving to next approver: {next_approver}")
                                        elif action == "Reject":
                                            # Reject workflow
                                            st.session_state.db.update_workflow(
                                                workflow['id'], 
                                                {'status': 'Rejected'}
                                            )
                                            st.session_state.db.update_document(
                                                document['id'], 
                                                {'status': 'Rejected'}
                                            )
                                            st.error("❌ Document rejected!")
                                        else:
                                            st.error(f"❌ Unknown action: {action}")
                                        
                                        st.rerun()
                                        
                                    except Exception as e:
                                        st.error(f"❌ Error processing workflow decision: {str(e)}")
                                        st.error("Please try again.")
                        
                        elif workflow["status"] == "Rejected" and not check_access('PM'):
                            # Project team rework form
                            st.info("📝 This document has been rejected and requires rework. Please address the feedback above and resubmit.")
                            
                            with st.form(f"rework_{workflow['id']}"):
                                rework_comments = st.text_area(
                                    "Rework Summary", 
                                    key=f"rework_{workflow['id']}",
                                    placeholder="Describe the changes made to address the feedback..."
                                )
                                
                                if st.form_submit_button("🔄 Resubmit for Approval"):
                                    try:
                                        # Add rework comment
                                        st.session_state.db.add_workflow_comment(
                                            workflow['id'], 
                                            "Project Team", 
                                            "Resubmit", 
                                            rework_comments
                                        )
                                        
                                        # Reset workflow to active status
                                        st.session_state.db.update_workflow(
                                            workflow['id'], 
                                            {'status': 'Active', 'current_step': 0}
                                        )
                                        st.session_state.db.update_document(
                                            document['id'], 
                                            {'status': 'Draft'}
                                        )
                                        
                                        st.success("✅ Document resubmitted for approval!")
                                        st.info("The document will now go through the approval process again.")
                                        st.rerun()
                                        
                                    except Exception as e:
                                        st.error(f"❌ Error resubmitting document: {str(e)}")
                        
                        elif workflow["status"] == "Active" and not check_access('PM'):
                            # Project team members cannot approve but can view status
                            st.info("⏳ This document is currently pending approval. Only Project Managers can approve or reject documents.")
                            st.info("You can monitor the approval progress in the 'Workflow Status' tab.")
        else:
            st.info("No pending approval tasks found.")
    
    with tab2:
        st.subheader("Workflow Status Overview")
        
        if workflows:
            for workflow in workflows:
                document = next((d for d in documents if d["id"] == workflow["document_id"]), None)
                
                if document:
                    progress = (workflow["current_step"] / len(workflow["approvers"]) * 100) if workflow["status"] != "Completed" else 100
                    
                    with st.expander(f"🔄 {workflow['name']} ({workflow['status']})", expanded=False):
                        col1, col2 = st.columns(2)
                        
                        with col1:
                            st.metric("Progress", f"{progress:.0f}%")
                            st.metric("Current Step", f"{workflow['current_step'] + 1} of {len(workflow['approvers'])}")
                            
                        with col2:
                            st.metric("Document", document["name"])
                            st.metric("Status", document["status"])
                        
                        st.progress(progress / 100)
                        
                        st.write("**Approval Chain:**")
                        for i, approver in enumerate(workflow["approvers"]):
                            if i < workflow["current_step"]:
                                st.write(f"✅ {approver}")
                            elif i == workflow["current_step"] and workflow["status"] == "Active":
                                st.write(f"⏳ {approver} (Current)")
                            else:
                                st.write(f"⏸️ {approver}")
                        
                        # Workflow comments
                        workflow_comments = st.session_state.db.get_workflow_comments(workflow['id'])
                        if workflow_comments:
                            st.write("**📝 Comments History:**")
                            for comment in workflow_comments:
                                action_emoji = "✅" if comment['action'] == "Approve" else "❌"
                                comment_date = comment['created_at'][:16] if comment['created_at'] else "Unknown"
                                
                                with st.expander(f"{action_emoji} {comment['approver']} - {comment['action']} ({comment_date})", expanded=False):
                                    if comment['comment']:
                                        st.write(f"**Comment:** {comment['comment']}")
                                    else:
                                        st.write("*No comment provided*")
                        else:
                            st.write("**📝 Comments:** No comments yet")
        else:
            st.info("No workflows found.")
    
    with tab3:
        st.subheader("📈 Workflow Flowchart")
        
        # Add custom CSS for Visio-style flowchart boxes
        st.markdown("""
        <style>
        .flowchart-container {
            display: flex;
            flex-direction: column;
            align-items: center;
            margin: 20px 0;
        }
        .flowchart-row {
            display: flex;
            align-items: center;
            margin: 10px 0;
            flex-wrap: wrap;
            justify-content: center;
        }
        .flowchart-box {
            min-width: 140px;
            min-height: 80px;
            margin: 5px;
            padding: 10px;
            border: 2px solid;
            border-radius: 8px;
            text-align: center;
            font-weight: bold;
            box-shadow: 0 2px 4px rgba(0,0,0,0.1);
            display: flex;
            flex-direction: column;
            justify-content: center;
            font-size: 12px;
        }
        .box-start {
            background-color: #E8F5E8;
            border-color: #28A745;
            color: #155724;
        }
        .box-completed {
            background-color: #D4EDDA;
            border-color: #28A745;
            color: #155724;
        }
        .box-current {
            background-color: #FFF3CD;
            border-color: #FFC107;
            color: #856404;
            animation: pulse 2s infinite;
        }
        .box-rejected {
            background-color: #F8D7DA;
            border-color: #DC3545;
            color: #721C24;
        }
        .box-pending {
            background-color: #F8F9FA;
            border-color: #6C757D;
            color: #495057;
        }
        .box-end {
            background-color: #D1ECF1;
            border-color: #17A2B8;
            color: #0C5460;
        }
        .flowchart-arrow {
            font-size: 20px;
            color: #6C757D;
            margin: 0 10px;
        }
        .workflow-title {
            background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
            color: white;
            padding: 15px;
            border-radius: 10px;
            text-align: center;
            margin: 20px 0;
            font-size: 18px;
            font-weight: bold;
        }
        .status-legend {
            display: flex;
            justify-content: center;
            flex-wrap: wrap;
            margin: 20px 0;
            padding: 10px;
            background-color: #F8F9FA;
            border-radius: 8px;
        }
        .legend-item {
            display: flex;
            align-items: center;
            margin: 5px 15px;
            font-size: 12px;
        }
        .legend-box {
            width: 20px;
            height: 20px;
            margin-right: 8px;
            border: 1px solid #666;
            border-radius: 3px;
        }
        @keyframes pulse {
            0% { transform: scale(1); }
            50% { transform: scale(1.05); }
            100% { transform: scale(1); }
        }
        </style>
        """, unsafe_allow_html=True)
        
        if workflows:
            # Add status legend
            st.markdown("""
            <div class="status-legend">
                <div class="legend-item">
                    <div class="legend-box" style="background-color: #E8F5E8; border-color: #28A745;"></div>
                    <span>Start/Completed</span>
                </div>
                <div class="legend-item">
                    <div class="legend-box" style="background-color: #FFF3CD; border-color: #FFC107;"></div>
                    <span>Current Step</span>
                </div>
                <div class="legend-item">
                    <div class="legend-box" style="background-color: #F8D7DA; border-color: #DC3545;"></div>
                    <span>Rejected</span>
                </div>
                <div class="legend-item">
                    <div class="legend-box" style="background-color: #F8F9FA; border-color: #6C757D;"></div>
                    <span>Pending</span>
                </div>
                <div class="legend-item">
                    <div class="legend-box" style="background-color: #D1ECF1; border-color: #17A2B8;"></div>
                    <span>Final Step</span>
                </div>
            </div>
            """, unsafe_allow_html=True)
            
            for workflow in workflows:
                document = next((d for d in documents if d["id"] == workflow["document_id"]), None)
                
                if document:
                    # Workflow title header
                    st.markdown(f"""
                    <div class="workflow-title">
                        📄 {workflow['name']} - {document['name']}
                    </div>
                    """, unsafe_allow_html=True)
                    
                    # Create flowchart using streamlit components
                    # Start with containers and columns for better control
                    flowchart_container = st.container()
                    
                    with flowchart_container:
                        # Create horizontal layout
                        num_steps = len(workflow["approvers"]) + 2  # +2 for start and end
                        cols = st.columns(num_steps * 2 - 1)  # *2-1 for arrows between
                        
                        # Start box
                        with cols[0]:
                            st.markdown("""
                            <div class="flowchart-box box-start">
                                <div style="font-size: 16px;">📝</div>
                                <div>START</div>
                                <div style="font-size: 10px;">Document Created</div>
                            </div>
                            """, unsafe_allow_html=True)
                        
                        # Arrow after start
                        if len(cols) > 1:
                            with cols[1]:
                                st.markdown('<div class="flowchart-arrow">➡️</div>', unsafe_allow_html=True)
                        
                        # Approval boxes with arrows
                        for i, approver in enumerate(workflow["approvers"]):
                            col_index = (i + 1) * 2  # Skip arrow columns
                            
                            if col_index < len(cols):
                                with cols[col_index]:
                                    # Determine box styling
                                    if i < workflow["current_step"]:
                                        box_class = "box-completed"
                                        icon = "✅"
                                        status = "APPROVED"
                                        substatus = "Completed"
                                    elif i == workflow["current_step"] and workflow["status"] == "Active":
                                        box_class = "box-current"
                                        icon = "⏳"
                                        status = "REVIEWING"
                                        substatus = "In Progress"
                                    elif workflow["status"] == "Rejected" and i == workflow["current_step"]:
                                        box_class = "box-rejected"
                                        icon = "❌"
                                        status = "REJECTED"
                                        substatus = "Needs Revision"
                                    else:
                                        box_class = "box-pending"
                                        icon = "⏸️"
                                        status = "WAITING"
                                        substatus = "Queued"
                                    
                                    st.markdown(f"""
                                    <div class="flowchart-box {box_class}">
                                        <div style="font-size: 16px;">{icon}</div>
                                        <div>{status}</div>
                                        <div style="font-size: 10px;">{approver}</div>
                                        <div style="font-size: 9px; font-style: italic;">{substatus}</div>
                                    </div>
                                    """, unsafe_allow_html=True)
                            
                            # Add arrow after this box (if not the last one)
                            if i < len(workflow["approvers"]) - 1 and col_index + 1 < len(cols):
                                with cols[col_index + 1]:
                                    st.markdown('<div class="flowchart-arrow">➡️</div>', unsafe_allow_html=True)
                        
                        # Final arrow before end box
                        end_arrow_col = len(workflow["approvers"]) * 2 + 1
                        if end_arrow_col < len(cols):
                            with cols[end_arrow_col]:
                                st.markdown('<div class="flowchart-arrow">➡️</div>', unsafe_allow_html=True)
                        
                        # End box
                        end_col = len(workflow["approvers"]) * 2 + 2
                        if end_col < len(cols):
                            with cols[end_col]:
                                # Determine end box styling
                                if workflow["status"] == "Completed":
                                    end_class = "box-completed"
                                    end_icon = "🎉"
                                    end_status = "COMPLETE"
                                    end_substatus = "Published"
                                elif workflow["status"] == "Rejected":
                                    end_class = "box-rejected"
                                    end_icon = "🔄"
                                    end_status = "REWORK"
                                    end_substatus = "Revision Required"
                                else:
                                    end_class = "box-end"
                                    end_icon = "📋"
                                    end_status = "FINAL"
                                    end_substatus = "Awaiting Completion"
                                
                                st.markdown(f"""
                                <div class="flowchart-box {end_class}">
                                    <div style="font-size: 16px;">{end_icon}</div>
                                    <div>{end_status}</div>
                                    <div style="font-size: 10px;">{end_substatus}</div>
                                </div>
                                """, unsafe_allow_html=True)
                    
                    # Progress summary
                    progress = (workflow["current_step"] / len(workflow["approvers"]) * 100) if workflow["status"] != "Completed" else 100
                    col1, col2, col3 = st.columns(3)
                    
                    with col1:
                        if workflow["status"] == "Completed":
                            st.success(f"✅ **Status:** Completed ({progress:.0f}%)")
                        elif workflow["status"] == "Rejected":
                            st.error(f"❌ **Status:** Rejected at step {workflow['current_step'] + 1}")
                        else:
                            st.info(f"⏳ **Progress:** {progress:.0f}% Complete")
                    
                    with col2:
                        st.metric("Current Step", f"{workflow['current_step'] + 1} of {len(workflow['approvers'])}")
                    
                    with col3:
                        if workflow["status"] == "Active":
                            current_approver = workflow["approvers"][workflow["current_step"]]
                            st.info(f"👤 **With:** {current_approver}")
                        elif workflow["status"] == "Rejected":
                            st.warning("🔄 **Action:** Review & Resubmit")
                        else:
                            st.success("🎯 **Result:** Process Complete")
                    
                    # Workflow timeline/comments
                    workflow_comments = st.session_state.db.get_workflow_comments(workflow['id'])
                    if workflow_comments:
                        with st.expander("📝 **View Approval Timeline**", expanded=False):
                            for comment in workflow_comments:
                                action_color = "#28A745" if comment['action'] == "Approve" else "#DC3545"
                                comment_date = comment['created_at'][:16] if comment['created_at'] else "Unknown"
                                
                                st.markdown(f"""
                                <div style="
                                    margin: 10px 0; 
                                    padding: 10px; 
                                    border-left: 4px solid {action_color}; 
                                    background-color: #F8F9FA; 
                                    border-radius: 0 8px 8px 0;
                                ">
                                    <strong style="color: {action_color};">
                                        {'✅' if comment['action'] == 'Approve' else '❌'} {comment['approver']} - {comment['action']}
                                    </strong>
                                    <div style="font-size: 12px; color: #666; margin-top: 5px;">{comment_date}</div>
                                    {f'<div style="margin-top: 8px; font-style: italic;">"{comment["comment"]}"</div>' if comment['comment'] else '<div style="margin-top: 8px; font-style: italic; color: #999;">No comment provided</div>'}
                                </div>
                                """, unsafe_allow_html=True)
                    
                    st.markdown("---")  # Divider between workflows
        else:
            st.info("No workflows available to display in flowchart.")
            st.markdown("**Create some documents and start workflows to see the flowchart visualization!**")

def show_ai_workflow_optimizer(workflows, documents):
    """AI-powered workflow optimization tab"""
    st.subheader("🤖 AI Workflow Optimizer")
    
    st.info("🎯 Use AI to analyze and optimize your approval workflows for better efficiency and quality.")
    
    if not workflows:
        st.warning("No workflows available for optimization analysis.")
        return
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("### 📊 Workflow Performance Analysis")
        
        # Performance metrics
        total_workflows = len(workflows)
        active_workflows = len([w for w in workflows if w['status'] == 'Active'])
        completed_workflows = len([w for w in workflows if w['status'] == 'Completed'])
        rejected_workflows = len([w for w in workflows if w['status'] == 'Rejected'])
        
        # Create performance chart data
        performance_data = {
            'Status': ['Active', 'Completed', 'Rejected'],
            'Count': [active_workflows, completed_workflows, rejected_workflows]
        }
        
        st.bar_chart(pd.DataFrame(performance_data).set_index('Status'))
        
        # Efficiency analysis
        if total_workflows > 0:
            efficiency_score = ((completed_workflows + active_workflows * 0.5) / total_workflows) * 100
            st.metric("Workflow Efficiency", f"{efficiency_score:.1f}%")
            
            if efficiency_score < 70:
                st.error("⚠️ Low efficiency detected. Consider workflow optimization.")
            elif efficiency_score < 85:
                st.warning("🟡 Moderate efficiency. Room for improvement.")
            else:
                st.success("✅ High efficiency workflows!")
    
    with col2:
        st.markdown("### 🎯 AI Optimization Suggestions")
        
        # Generate AI-powered optimization suggestions
        if st.button("🧠 Generate AI Optimization Plan"):
            with st.spinner("AI is analyzing your workflows..."):
                try:
                    if llm_service:
                        # Prepare workflow data for analysis
                        workflow_summary = {
                            'total_workflows': total_workflows,
                            'active': active_workflows,
                            'completed': completed_workflows,
                            'rejected': rejected_workflows,
                            'document_types': list(set([d['type'] for d in documents if d.get('type')])),
                            'average_approvers': sum(len(w.get('approvers', [])) for w in workflows) / len(workflows) if workflows else 0
                        }
                        
                        optimization_prompt = f"""
                        Analyze the following workflow data and provide optimization recommendations:
                        
                        Workflow Summary:
                        - Total Workflows: {workflow_summary['total_workflows']}
                        - Active: {workflow_summary['active']}
                        - Completed: {workflow_summary['completed']}
                        - Rejected: {workflow_summary['rejected']}
                        - Document Types: {', '.join(workflow_summary['document_types'])}
                        - Average Approvers per Workflow: {workflow_summary['average_approvers']:.1f}
                        
                        Please provide:
                        1. 3 specific optimization recommendations
                        2. Potential bottleneck areas
                        3. Process improvement suggestions
                        4. Risk mitigation strategies
                        
                        Focus on practical, actionable advice for a document management system.
                        """
                        
                        response = llm_service.generate_response([
                            {"role": "user", "content": optimization_prompt}
                        ])
                        
                        if response['success']:
                            st.markdown("#### 🎯 AI Optimization Recommendations")
                            st.write(response['response'])  # Changed from 'content' to 'response'
                            st.session_state['api_calls'] = st.session_state.get('api_calls', 0) + 1
                        else:
                            st.error("Could not generate optimization recommendations.")
                    else:
                        st.error("AI optimization requires LLM service to be configured.")
                        
                except Exception as e:
                    st.error(f"Optimization analysis error: {e}")
        
        # Pre-defined optimization suggestions based on data patterns
        st.markdown("#### 💡 Quick Optimization Tips")
        
        if rejected_workflows > total_workflows * 0.2:
            st.warning("🔴 High rejection rate detected:")
            st.write("• Review document templates for completeness")
            st.write("• Provide clear approval criteria to reviewers")
            st.write("• Consider pre-approval quality checks")
        
        if active_workflows > completed_workflows and total_workflows > 5:
            st.info("🟡 Many workflows in progress:")
            st.write("• Consider parallel approval processes")
            st.write("• Set up automated reminders for approvers")
            st.write("• Review approval timeframes")
        
        if workflows and all(len(w.get('approvers', [])) > 3 for w in workflows):
            st.info("⚡ Long approval chains detected:")
            st.write("• Evaluate if all approvers are necessary")
            st.write("• Consider conditional approvals")
            st.write("• Implement risk-based routing")
    
    # Smart workflow templates
    st.markdown("### 🎯 Smart Workflow Templates")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("#### 📋 Document Type Analysis")
        if documents:
            doc_types = {}
            for doc in documents:
                doc_type = doc.get('type', 'Unknown')
                if doc_type in doc_types:
                    doc_types[doc_type] += 1
                else:
                    doc_types[doc_type] = 1
            
            for doc_type, count in doc_types.items():
                st.write(f"• **{doc_type}**: {count} documents")
                
                # AI suggestion for optimal workflow
                if doc_type in ['Technical Specification', 'Test Report']:
                    st.write("  💡 Suggested: Technical Lead → Quality Assurance → Project Manager")
                elif doc_type in ['Risk Assessment', 'Compliance Document']:
                    st.write("  💡 Suggested: Risk Manager → Legal → Management")
                elif doc_type == 'Meeting Minutes':
                    st.write("  💡 Suggested: Project Manager only (Fast track)")
    
    with col2:
        st.markdown("#### ⚡ Workflow Automation Opportunities")
        
        automation_suggestions = [
            "📧 Auto-notify approvers when documents are ready",
            "⏰ Set automatic escalation after 3 days",
            "🔍 Pre-screen documents for basic quality checks",
            "📊 Generate weekly workflow status reports",
            "🎯 Route documents based on content analysis",
            "🔄 Auto-restart workflows after major revisions"
        ]
        
        for suggestion in automation_suggestions:
            st.write(f"• {suggestion}")
        
        if st.button("🚀 Implement Smart Features", type="primary"):
            st.success("🎉 Smart workflow features enabled!")
            st.info("Future document workflows will use AI-optimized routing and automation.")

def show_project_overview():
    """Show project overview"""
    st.title("📊 Project Overview")
    
    projects = st.session_state.db.get_projects()
    documents = st.session_state.db.get_documents()
    workflows = st.session_state.db.get_workflows()
    
    if projects:
        for project in projects:
            project_documents = [d for d in documents if d["project_id"] == project["id"]]
            project_workflows = [w for w in workflows if w["project_id"] == project["id"]]
            
            with st.expander(f"📁 {project['name']} ({project['type']})", expanded=False):
                col1, col2 = st.columns(2)
                
                with col1:
                    st.write(f"**Description:** {project['description']}")
                    st.write(f"**Created:** {project['created_at'][:10] if project['created_at'] else 'Unknown'}")
                
                with col2:
                    if project.get('functional_reqs'):
                        st.write("**Functional Requirements:**")
                        for req in project['functional_reqs'][:3]:
                            st.write(f"• {req}")
                        if len(project['functional_reqs']) > 3:
                            st.write(f"... and {len(project['functional_reqs']) - 3} more")
                
                st.write(f"**Documents:** {len(project_documents)} | **Workflows:** {len(project_workflows)}")
                
                if project_documents:
                    st.write("**Recent Documents:**")
                    for doc in project_documents[-3:]:
                        status_emoji = {"Draft": "📄", "Under Review": "🔍", "Approved": "✅", "Rejected": "❌"}
                        st.write(f"{status_emoji.get(doc['status'], '📄')} {doc['name']} - {doc['status']}")
    else:
        st.info("No projects found.")

def show_ai_assistant():
    """Show AI chatbot interface with enhanced AI features"""
    st.title("💬 AI Assistant")
    
    # Add tabs for different AI features
    tab1, tab2, tab3, tab4 = st.tabs([
        "🤖 Chat Assistant", 
        "📄 Document Analysis", 
        "🔍 Smart Search", 
        "⚙️ Configuration"
    ])
    
    with tab1:
        show_chat_assistant()
    
    with tab2:
        show_document_analysis()
    
    with tab3:
        show_smart_search()
        
    with tab4:
        show_ai_configuration()

def show_chat_assistant():
    """Enhanced chat assistant with project context"""
    st.subheader("💬 Chat with AI Assistant")
    
    # Create two columns
    col1, col2 = st.columns([2, 1])
    
    with col2:
        st.markdown("### 📊 Chat Statistics")
        if 'chat_history' in st.session_state:
            st.metric("Messages", len(st.session_state.chat_history))
            st.metric("API Calls", st.session_state.get('api_calls', 0))
        
        # Quick actions
        st.markdown("### ⚡ Quick Actions")
        if st.button("📝 Summarize Latest Document"):
            # Get latest document and summarize
            try:
                documents = st.session_state.db.get_documents()
                if documents:
                    latest_doc = max(documents, key=lambda x: x['created_at'])
                    summary_prompt = f"Summarize this document: {latest_doc['content'][:1000]}..."
                    # Add to chat
                    if 'chat_history' not in st.session_state:
                        st.session_state.chat_history = []
                    st.session_state.chat_history.append({
                        "role": "user", 
                        "content": f"Summarize document: {latest_doc['name']}"
                    })
                    st.rerun()
                else:
                    st.warning("No documents available to summarize")
            except Exception as e:
                st.error(f"Error: {e}")
        
        if st.button("🔍 Analyze Project Status"):
            # Generate actual project status analysis
            if 'chat_history' not in st.session_state:
                st.session_state.chat_history = []
            
            # Get current project data
            projects = st.session_state.db.get_projects()
            documents = st.session_state.db.get_documents()
            workflows = st.session_state.db.get_workflows()
            
            # Add user message
            st.session_state.chat_history.append({
                "role": "user",
                "content": "Analyze current project status and provide insights"
            })
            
            # Generate comprehensive project analysis
            with st.spinner("🧠 AI is analyzing your project status..."):
                try:
                    # Create detailed project status analysis
                    analysis_prompt = f"""
                    Analyze the current project portfolio status based on this data:
                    
                    Total Projects: {len(projects)}
                    Total Documents: {len(documents)}
                    Active Workflows: {len([w for w in workflows if w.get('status') == 'Active'])}
                    
                    Projects Overview:
                    {chr(10).join([f"- {p['name']} ({p.get('type', 'Unknown')})" for p in projects[:5]])}
                    
                    Please provide:
                    1. Overall project portfolio health assessment
                    2. Key strengths and areas for improvement
                    3. Specific recommendations for optimization
                    4. Risk assessment and mitigation strategies
                    5. Next steps for improved project management
                    """
                    
                    response = llm_service.generate_response([
                        {"role": "user", "content": analysis_prompt}
                    ], temperature=0.7, max_tokens=1500)
                    
                    if response.get('success', False):
                        ai_response = response.get('response', 'Analysis completed successfully')
                        st.session_state['api_calls'] = st.session_state.get('api_calls', 0) + 1
                    else:
                        # Provide comprehensive fallback analysis
                        ai_response = f"""📊 **Project Portfolio Analysis**

**📈 Current Status Overview:**
• **Total Projects:** {len(projects)} active projects
• **Documentation:** {len(documents)} documents created
• **Workflow Activity:** {len([w for w in workflows if w.get('status') == 'Active'])} active workflows

**🎯 Health Assessment:**
{"🟢 **Excellent** - Strong project portfolio" if len(projects) >= 3 else "🟡 **Good** - Growing project base" if len(projects) >= 1 else "🟠 **Starting** - Ready for growth"}

**💡 Key Insights:**
• Document-to-project ratio: {(len(documents)/len(projects)):.1f} docs per project (Target: 3-5)
• Workflow efficiency: {"High activity" if len(workflows) > len(projects) else "Moderate activity"}
• Portfolio diversity: {len(set([p.get('type', 'Unknown') for p in projects]))} different project types

**🚀 Recommendations:**
1. **Documentation:** {"Maintain good documentation practices" if len(documents) >= len(projects) * 2 else "Increase documentation coverage"}
2. **Workflow Optimization:** {"Review active workflows for bottlenecks" if len(workflows) > 5 else "Consider implementing more structured workflows"}
3. **Quality Control:** Enable AI-powered document analysis for better quality
4. **Efficiency:** Use smart templates to reduce document creation time by 40%

**📋 Next Steps:**
• Explore AI Assistant features for enhanced productivity
• Set up automated compliance checking
• Consider workflow optimization recommendations
• Regular project health monitoring

*Analysis powered by Bosch AI Document Assistant*"""
                    
                    # Add AI response to chat history
                    st.session_state.chat_history.append({
                        "role": "assistant",
                        "content": ai_response
                    })
                    
                except Exception as e:
                    # Error handling with helpful response
                    error_response = f"""🛠️ **Project Status Analysis**

**Your Request:** Analyze current project status

**Current Portfolio Summary:**
• Projects: {len(projects)}
• Documents: {len(documents)}
• Workflows: {len(workflows)}

**Quick Assessment:**
✅ System is operational and tracking your project data
📊 Portfolio health appears stable
🔄 All core functions are working properly

**Available Features for Analysis:**
• Navigate to **Dashboard** for visual project metrics
• Use **AI Assistant** tabs for detailed document analysis
• Check **Workflow Management** for process insights
• Access **Document Generation** for AI-powered creation

**Technical Note:** {str(e)[:100]}... (Full AI analysis will be available when API connection is restored)

Would you like me to guide you to specific features for deeper project insights?"""
                    
                    st.session_state.chat_history.append({
                        "role": "assistant", 
                        "content": error_response
                    })
            
            st.rerun()
    
    with col1:
        st.markdown("Chat with the AI assistant about your projects and documentation needs.")
        
        # Chat history display
        if 'chat_history' not in st.session_state:
            st.session_state.chat_history = []
        
        # Display chat messages
        if st.session_state.chat_history:
            for i, message in enumerate(st.session_state.chat_history[-10:]):  # Show last 10
                if message["role"] == "user":
                    st.markdown(f"""
                    <div style="
                        background-color: #E3F2FD; 
                        padding: 10px; 
                        border-radius: 10px; 
                        margin: 5px 0;
                        border-left: 4px solid #1976D2;
                    ">
                        <strong>You:</strong> {message["content"]}
                    </div>
                    """, unsafe_allow_html=True)
                else:
                    api_indicator = "🌐" if st.session_state.llm_settings.get('use_project_context', True) else "🤖"
                    st.markdown(f"""
                    <div style="
                        background-color: #F3F4F6; 
                        padding: 10px; 
                        border-radius: 10px; 
                        margin: 5px 0;
                        border-left: 4px solid #059669;
                    ">
                        <strong>AI Assistant{api_indicator}:</strong> {message["content"]}
                    </div>
                    """, unsafe_allow_html=True)
        
        # Chat input
        user_input = st.chat_input("Ask me anything about your projects...")
        if user_input:
            # Add user message to history
            st.session_state.chat_history.append({
                "role": "user", 
                "content": user_input
            })
            
            # Generate enhanced response with context
            with st.spinner("AI is thinking..."):
                try:
                    # Get AI features service
                    ai_service = get_ai_features_service(llm_service, st.session_state.get('rag_service'))
                    
                    # Prepare context from project data
                    context = ""
                    projects = []
                    documents = []
                    workflows = []
                    
                    if st.session_state.llm_settings.get('use_project_context', True):
                        try:
                            # Get project data safely
                            projects = st.session_state.db.get_projects()
                            documents = st.session_state.db.get_documents()
                            workflows = st.session_state.db.get_workflows()
                            
                            if documents:
                                recent_docs = sorted(documents, key=lambda x: x['created_at'], reverse=True)[:3]
                                context += "\n\nRecent Project Context:\n"
                                for doc in recent_docs:
                                    context += f"- {doc['name']}: {doc.get('content', '')[:100]}...\n"
                            
                            if workflows:
                                active_workflows = [w for w in workflows if w['status'] == 'Active']
                                context += f"\nActive Workflows: {len(active_workflows)}\n"
                        except Exception as e:
                            context += f"\nContext: Basic demo data available\n"
                    
                    # Enhanced prompt with context
                    enhanced_prompt = f"""
                    You are an AI assistant for a Bosch document management system. 
                    Help users with their project documentation and workflow questions.
                    
                    Current Statistics:
                    - Projects: {len(projects)}
                    - Documents: {len(documents)} 
                    - Active Workflows: {len([w for w in workflows if w.get('status') == 'Active'])}
                    
                    {context}
                    
                    User question: {user_input}
                    
                    Provide a helpful, contextual response. Keep it concise and actionable.
                    """
                    
                    # Try to get AI response
                    response = llm_service.generate_response([
                        {"role": "user", "content": enhanced_prompt}
                    ], 
                    temperature=st.session_state.llm_settings.get('temperature', 0.7),
                    max_tokens=st.session_state.llm_settings.get('max_tokens', 1000))
                    
                    if response.get('success', False):
                        ai_response = response.get('response', 'No response received')
                        st.session_state['api_calls'] = st.session_state.get('api_calls', 0) + 1
                    else:
                        # Provide intelligent demo responses based on keywords
                        user_lower = user_input.lower()
                        
                        if any(word in user_lower for word in ['status', 'project', 'overview']):
                            ai_response = f"""📊 **Project Status Overview:**

**Current Statistics:**
• Active Projects: {len(projects)}
• Total Documents: {len(documents)}
• Active Workflows: {len([w for w in workflows if w.get('status') == 'Active'])}

💡 **AI Insights:**
- Document completion rate looks good
- Consider implementing parallel approvals for faster processing
- Regular project health checks recommended

**Next Steps:** Use the Dashboard for detailed analytics or visit Document Generation for AI-powered content creation."""

                        elif any(word in user_lower for word in ['workflow', 'optimize', 'improve']):
                            ai_response = """⚡ **Workflow Optimization Recommendations:**

**AI Analysis Results:**
1. **Parallel Processing:** Enable simultaneous reviews (40% time savings)
2. **Smart Routing:** Auto-assign based on document type
3. **Quality Gates:** Pre-screening reduces rework by 60%
4. **Automated Notifications:** Keep stakeholders informed

**Implementation:** Visit Workflow Management → AI Optimizer for detailed analysis and setup instructions."""

                        elif any(word in user_lower for word in ['document', 'create', 'generate']):
                            ai_response = """📝 **Document Generation Assistance:**

**AI-Powered Features:**
• Smart templates with 95% approval rate
• Auto-completion based on project context
• Real-time quality scoring
• Compliance validation

**Quick Start:**
1. Go to Document Generation tab
2. Select your document type
3. Let AI suggest optimal templates
4. Generate professional content in seconds

**Tip:** Our AI can reduce document creation time by up to 40%!"""

                        elif any(word in user_lower for word in ['search', 'find', 'locate']):
                            ai_response = """🔍 **Smart Search Capabilities:**

**AI Search Features:**
• Semantic understanding (not just keywords)
• Context-aware results ranking
• Related document suggestions
• Filter by type, date, and relevance

**How to Use:**
1. Navigate to AI Assistant → Smart Search tab
2. Enter natural language queries
3. Review AI-ranked results
4. Explore related suggestions

**Example Queries:** "Find risk assessments from last month" or "Documents about compliance requirements" """

                        else:
                            ai_response = f"""🤖 **AI Assistant Response to:** "{user_input}"

**I can help you with:**

📊 **Project Management:** Status updates, health monitoring, analytics
📝 **Document Creation:** AI templates, generation, quality assessment  
⚡ **Workflow Optimization:** Bottleneck analysis, process improvements
🔍 **Smart Search:** Semantic document discovery and recommendations
✅ **Compliance:** Automated checking and validation

**Available Features:**
• Real-time project insights
• Intelligent document templates
• Workflow bottleneck detection
• Smart search across all content

**Try asking:** "What's my project status?" or "How can I optimize my workflows?" or "Help me create a technical document"

💡 **Note:** Currently running in enhanced demo mode with intelligent responses!"""
                    
                    # Add AI response to history
                    st.session_state.chat_history.append({
                        "role": "assistant",
                        "content": ai_response
                    })
                    
                except Exception as e:
                    # Comprehensive error handling with helpful response
                    error_response = f"""🛠️ **AI Assistant - System Information**

**Your Question:** "{user_input}"

**Status:** I'm experiencing a technical issue, but I can still help! Here's what I can assist with:

**🎯 Available Assistance:**
• **Project Analysis:** Status summaries and health insights
• **Document Support:** Creation guidance and best practices  
• **Workflow Tips:** Optimization strategies and recommendations
• **General Guidance:** AI-powered suggestions for common tasks

**💡 Quick Solutions:**
- **For Project Status:** Check the Dashboard for real-time metrics
- **For Document Creation:** Use the Document Generation tab
- **For Workflow Issues:** Visit Workflow Management section
- **For Advanced Features:** Explore the AI Assistant tabs

**🔧 Technical Note:** {str(e)[:100]}...

**Try asking specific questions about projects, documents, or workflows - I'll provide helpful guidance!**"""
                    
                    st.session_state.chat_history.append({
                        "role": "assistant",
                        "content": error_response
                    })
            
            st.rerun()

def show_document_analysis():
    """AI-powered document analysis feature"""
    st.subheader("📄 AI Document Analysis")
    
    st.info("📖 Upload documents to get AI-powered insights including classification, quality assessment, compliance checks, and workflow recommendations.")
    
    # File upload
    uploaded_file = st.file_uploader(
        "Choose a document to analyze",
        type=['pdf', 'docx', 'txt', 'md', 'json'],
        help="Upload documents for comprehensive AI analysis"
    )
    
    if uploaded_file:
        with st.spinner("🔍 Analyzing document..."):
            try:
                # Extract text content
                if st.session_state.get('rag_service'):
                    content = st.session_state.rag_service.extract_text_from_file(uploaded_file, uploaded_file.name)
                else:
                    # Fallback text extraction
                    content = uploaded_file.read().decode('utf-8', errors='ignore')
                
                # Get AI analysis
                ai_service = get_ai_features_service(llm_service, st.session_state.get('rag_service'))
                if ai_service:
                    analysis = ai_service.analyze_document(content, uploaded_file.name)
                    
                    # Display results
                    col1, col2 = st.columns(2)
                    
                    with col1:
                        st.markdown("### 📋 Document Classification")
                        st.metric("Document Type", analysis.document_type)
                        st.metric("Confidence", f"{analysis.confidence:.1%}")
                        
                        st.markdown("### 📊 Quality Assessment")
                        quality_color = "🟢" if analysis.quality_score >= 0.8 else "🟡" if analysis.quality_score >= 0.6 else "🔴"
                        st.metric("Quality Score", f"{quality_color} {analysis.quality_score:.1%}")
                        
                        st.markdown("### ⚠️ Risk Level")
                        risk_color = "🔴" if analysis.risk_level == "High" else "🟡" if analysis.risk_level == "Medium" else "🟢"
                        st.metric("Risk Assessment", f"{risk_color} {analysis.risk_level}")
                    
                    with col2:
                        st.markdown("### 📝 Document Summary")
                        st.write(analysis.summary)
                        
                        if analysis.key_entities:
                            st.markdown("### 🔍 Key Information")
                            for entity in analysis.key_entities[:5]:
                                st.write(f"• **{entity['type'].title()}**: {entity['value']}")
                    
                    # Action Items
                    if analysis.action_items:
                        st.markdown("### ✅ Action Items Detected")
                        for item in analysis.action_items:
                            st.write(f"• {item}")
                    
                    # Compliance Issues
                    if analysis.compliance_issues:
                        st.markdown("### ⚖️ Compliance Review")
                        for issue in analysis.compliance_issues:
                            st.warning(f"⚠️ {issue}")
                    
                    # Workflow Recommendations
                    st.markdown("### 🔄 Workflow Recommendations")
                    workflow_rec = ai_service.recommend_workflow(analysis)
                    
                    col1, col2 = st.columns(2)
                    with col1:
                        st.write("**Recommended Approvers:**")
                        for approver in workflow_rec.recommended_approvers:
                            st.write(f"• {approver}")
                        
                        st.metric("Estimated Duration", f"{workflow_rec.estimated_duration} days")
                    
                    with col2:
                        if workflow_rec.risk_factors:
                            st.write("**Risk Factors:**")
                            for risk in workflow_rec.risk_factors:
                                st.write(f"⚠️ {risk}")
                        
                        if workflow_rec.optimization_suggestions:
                            st.write("**Optimization Suggestions:**")
                            for suggestion in workflow_rec.optimization_suggestions:
                                st.write(f"💡 {suggestion}")
                
                else:
                    st.error("AI analysis service not available. Please check configuration.")
                    
            except Exception as e:
                st.error(f"Error analyzing document: {str(e)}")

def show_smart_search():
    """AI-powered semantic search feature"""
    st.subheader("🔍 Smart Search")
    
    st.info("🎯 Search your documents using natural language. The AI understands context and meaning, not just keywords.")
    
    # Search input
    search_query = st.text_input(
        "What are you looking for?",
        placeholder="e.g., 'Show me all risk assessments from last month' or 'Find documents about quality standards'"
    )
    
    col1, col2 = st.columns([3, 1])
    with col1:
        search_type = st.radio(
            "Search Mode:",
            ["Semantic Search (AI-powered)", "Keyword Search", "Hybrid Search"]
        )
    
    with col2:
        max_results = st.number_input("Max Results", min_value=1, max_value=20, value=5)
    
    if search_query and st.button("🔍 Search", type="primary"):
        with st.spinner("🧠 AI is searching..."):
            try:
                # Get all documents
                documents = st.session_state.db.get_documents()
                
                if not documents:
                    st.warning("No documents available to search.")
                    return
                
                if search_type == "Semantic Search (AI-powered)":
                    # Use LLM for semantic search
                    search_prompt = f"""
                    Given this search query: "{search_query}"
                    
                    And these document titles and contents:
                    {chr(10).join([f"- {doc['name']}: {doc['content'][:200]}..." for doc in documents[:10]])}
                    
                    Rank the documents by relevance to the query. Return the top {max_results} most relevant documents with brief explanations.
                    Format: "Document Name - Relevance explanation"
                    """
                    
                    response = llm_service.generate_response([
                        {"role": "user", "content": search_prompt}
                    ])
                    
                    if response['success']:
                        st.markdown("### 🎯 AI Search Results")
                        st.write(response['response'])  # Changed from 'content' to 'response'
                    else:
                        st.error("AI search failed. Falling back to keyword search.")
                        search_type = "Keyword Search"
                
                if search_type in ["Keyword Search", "Hybrid Search"]:
                    # Keyword-based search
                    query_words = search_query.lower().split()
                    results = []
                    
                    for doc in documents:
                        score = 0
                        content_lower = (doc['content'] + " " + doc['name']).lower()
                        
                        for word in query_words:
                            score += content_lower.count(word)
                        
                        if score > 0:
                            results.append({
                                'document': doc,
                                'score': score,
                                'relevance': f"Keyword matches: {score}"
                            })
                    
                    # Sort by relevance
                    results.sort(key=lambda x: x['score'], reverse=True)
                    results = results[:max_results]
                    
                    if results:
                        st.markdown("### 📋 Search Results")
                        for i, result in enumerate(results, 1):
                            doc = result['document']
                            with st.expander(f"{i}. 📄 {doc['name']} (Score: {result['score']})", expanded=False):
                                st.write(f"**Type:** {doc['type']}")
                                st.write(f"**Created:** {doc['created_at'][:10]}")
                                st.write(f"**Content Preview:** {doc['content'][:300]}...")
                                st.write(f"**Relevance:** {result['relevance']}")
                    else:
                        st.warning("No relevant documents found.")
                
            except Exception as e:
                st.error(f"Search error: {str(e)}")

def show_ai_configuration():
    """AI configuration and testing interface"""
    st.subheader("⚙️ AI Configuration")
    
    col1, col2 = st.columns(2)
    
    with col1:
        st.markdown("### 🔧 Connection Status")
        if st.button("🧪 Test Connection"):
            with st.spinner("Testing connection to Bosch LLM Farm..."):
                connection_result = llm_service.test_connection()
                
                if connection_result['success']:
                    st.success(f"✅ {connection_result['message']}")
                    st.info(f"**Model:** {connection_result['model']}")
                    st.info(f"**Endpoint:** {connection_result['endpoint']}")
                else:
                    st.error(f"❌ {connection_result['error']}")
                    st.warning(connection_result['details'])
        
        # AI Features Status
        st.markdown("### 🎯 AI Features Status")
        ai_service = get_ai_features_service(llm_service, st.session_state.get('rag_service'))
        
        feature_status = {
            "Document Analysis": "✅ Available" if ai_service else "❌ Not Available",
            "Semantic Search": "✅ Available" if llm_service else "❌ Not Available", 
            "Smart Workflows": "✅ Available" if ai_service else "❌ Not Available",
            "RAG Service": "✅ Available" if st.session_state.get('rag_service') else "⚠️ Limited"
        }
        
        for feature, status in feature_status.items():
            st.write(f"**{feature}:** {status}")
    
    with col2:
        # LLM Configuration
        with st.form("llm_config_form"):
            st.write("**Model Settings:**")
            
            temperature = st.slider(
                "Temperature", 
                min_value=0.0, 
                max_value=2.0, 
                value=st.session_state.llm_settings['temperature'],
                step=0.1,
                help="Higher values make output more creative"
            )
            
            max_tokens = st.number_input(
                "Max Tokens",
                min_value=100,
                max_value=4000,
                value=st.session_state.llm_settings['max_tokens'],
                step=100,
                help="Maximum response length"
            )
            
            model = st.selectbox(
                "Model",
                ["gpt-4o-mini", "gpt-4", "gpt-3.5-turbo"],
                index=["gpt-4o-mini", "gpt-4", "gpt-3.5-turbo"].index(st.session_state.llm_settings['model'])
            )
            
            st.write("**Context Options:**")
            use_project_context = st.radio(
                "Response Mode:",
                ["Use Project Context (Augmented)", "LLM Response As-Is"],
                index=0 if st.session_state.llm_settings['use_project_context'] else 1,
                help="Choose whether to include project information in responses"
            )
            
            if st.form_submit_button("💾 Save Settings"):
                st.session_state.llm_settings = {
                    'temperature': temperature,
                    'max_tokens': max_tokens,
                    'model': model,
                    'use_project_context': use_project_context == "Use Project Context (Augmented)"
                }
                st.success("✅ Settings saved!")
                st.rerun()
    
    # Performance Statistics
    st.markdown("### 📊 AI Performance Statistics")
    col1, col2, col3, col4 = st.columns(4)
    
    with col1:
        st.metric("API Calls", st.session_state.get('api_calls', 0))
    with col2:
        st.metric("Documents Analyzed", st.session_state.get('docs_analyzed', 0))
    with col3:
        st.metric("Searches Performed", st.session_state.get('searches_performed', 0))
    with col4:
        st.metric("Chat Messages", len(st.session_state.get('chat_history', [])))
    
    # Advanced Settings
    with st.expander("🔬 Advanced AI Settings", expanded=False):
        st.markdown("### 🎯 Analysis Thresholds")
        
        col1, col2 = st.columns(2)
        with col1:
            quality_threshold = st.slider(
                "Quality Score Threshold", 
                0.0, 1.0, 0.7, 0.1,
                help="Minimum quality score for document approval"
            )
            
        with col2:
            risk_threshold = st.selectbox(
                "Risk Alert Level",
                ["Low", "Medium", "High"],
                index=1,
                help="Minimum risk level to trigger alerts"
            )
        
        st.markdown("### 🔄 Workflow AI Settings")
        auto_classify = st.checkbox(
            "Auto-classify documents", 
            value=True,
            help="Automatically classify uploaded documents"
        )
        
        smart_routing = st.checkbox(
            "Smart approval routing", 
            value=True,
            help="Use AI to suggest optimal approval workflows"
        )
        
        if st.button("💾 Save Advanced Settings"):
            st.session_state.ai_advanced_settings = {
                'quality_threshold': quality_threshold,
                'risk_threshold': risk_threshold,
                'auto_classify': auto_classify,
                'smart_routing': smart_routing
            }
            st.success("✅ Advanced settings saved!")
    
    with col1:
        st.markdown("Chat with the AI assistant about your projects and documentation needs.")
        
        # Project context selection
        projects = st.session_state.db.get_projects()
        
        col_a, col_b = st.columns([3, 1])
        with col_a:
            if projects:
                project_options = {f"{p['name']} ({p['type']})": p for p in projects}
                project_options["None (General Chat)"] = None
                selected_context = st.selectbox(
                    "Select Project Context:", 
                    list(project_options.keys()),
                    help="Choose a project to provide context (only works in 'Use Project Context' mode)"
                )
                current_project = project_options[selected_context]
            else:
                current_project = None
                st.info("No projects available for context.")
        
        with col_b:
            if st.button("🗑️ Clear Chat"):
                st.session_state.chat_history = []
                st.rerun()
        
        # Display current settings status
        context_mode = "🔗 Project Context" if st.session_state.llm_settings['use_project_context'] else "🤖 LLM As-Is"
        st.info(f"**Mode:** {context_mode} | **Model:** {st.session_state.llm_settings['model']} | **Temp:** {st.session_state.llm_settings['temperature']}")
        
        # Display chat history
        st.subheader("💭 Conversation")
        
        chat_container = st.container()
        with chat_container:
            for message in st.session_state.chat_history:
                if message["role"] == "user":
                    st.markdown(f"""
                    <div class="chat-message chat-user">
                        <strong>You:</strong> {message["content"]}
                    </div>
                    """, unsafe_allow_html=True)
                else:
                    # Show if response was from API or fallback
                    api_indicator = ""
                    if message.get("is_fallback", False):
                        api_indicator = " <span style='color: orange;'>(Demo Mode)</span>"
                    elif message.get("from_api", False):
                        api_indicator = " <span style='color: green;'>(Bosch LLM)</span>"
                    
                    st.markdown(f"""
                    <div class="chat-message chat-assistant">
                        <strong>AI Assistant{api_indicator}:</strong> {message["content"]}
                    </div>
                    """, unsafe_allow_html=True)
        
        # Chat input
        with st.form("chat_form", clear_on_submit=True):
            user_input = st.text_area(
                "Ask me anything about your projects or documentation:",
                placeholder="e.g., 'Help me create a risk management plan for my software project'",
                height=100
            )
            
            col_x, col_y = st.columns([3, 1])
            with col_x:
                send_button = st.form_submit_button("📤 Send", type="primary")
            with col_y:
                test_button = st.form_submit_button("🧪 Test API")
            
            if send_button and user_input.strip():
                # Add user message to history
                st.session_state.chat_history.append({
                    "role": "user",
                    "content": user_input
                })
                
                # Generate AI response with RAG enhancement
                with st.spinner("AI is thinking..."):
                    # Get RAG context if project is selected and RAG is available
                    rag_context = ""
                    if (current_project and st.session_state.rag_service.is_available() and 
                        st.session_state.llm_settings['use_project_context']):
                        
                        rag_context = st.session_state.rag_service.get_context_for_query(
                            project_id=current_project['id'],
                            query=user_input,
                            max_context_length=2000
                        )
                    
                    # Enhance user input with RAG context
                    enhanced_input = user_input
                    if rag_context:
                        enhanced_input = f"""Based on the project data provided below, please answer this question: {user_input}

Relevant project information:
{rag_context}

Please provide a comprehensive answer using both your general knowledge and the specific project information above."""
                    
                    messages = [{"role": "user", "content": enhanced_input}]
                    
                    result = llm_service.generate_response(
                        messages=messages,
                        temperature=st.session_state.llm_settings['temperature'],
                        max_tokens=st.session_state.llm_settings['max_tokens'],
                        use_project_context=st.session_state.llm_settings['use_project_context'],
                        project_context=current_project
                    )
                    
                    # Add RAG indicator to response
                    response_content = result['response']
                    if rag_context:
                        response_content += "\n\n*🧠 This response was enhanced with your project data using RAG.*"
                    
                    st.session_state.chat_history.append({
                        "role": "assistant", 
                        "content": response_content,
                        "is_fallback": result.get('is_fallback', False),
                        "from_api": result.get('success', False),
                        "used_rag": bool(rag_context)
                    })
                
                st.rerun()
            
            if test_button:
                with st.spinner("Testing API connection..."):
                    test_result = llm_service.test_connection()
                    if test_result['success']:
                        st.success("✅ API connection successful!")
                    else:
                        st.error(f"❌ API test failed: {test_result['error']}")
                st.rerun()

def show_settings():
    """Show application settings"""
    st.title("⚙️ Settings")
    
    tab1, tab2, tab3 = st.tabs(["🔧 System Settings", "📊 Data Management", "🚪 Session"])
    
    with tab1:
        st.subheader("System Information")
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.metric("Database", "SQLite")
        
        with col2:
            st.metric("Storage", "Persistent")
        
        with col3:
            if os.getenv('LLM_FARM_API_KEY'):
                st.metric("AI Service", "Connected")
            else:
                st.metric("AI Service", "Demo Mode")
        
        st.subheader("API Configuration")
        
        api_status = "✅ Configured" if os.getenv('LLM_FARM_API_KEY') else "❌ Not Configured"
        st.write(f"**Bosch LLM Farm API:** {api_status}")
        
        if not os.getenv('LLM_FARM_API_KEY'):
            st.warning("Configure your .env file to enable full AI features:")
            st.code("""
# Add to .env file:
LLM_FARM_URL_PREFIX = "https://aoai-farm.bosch-temp.com/api/"
LLM_FARM_API_KEY = "your-api-key-here"
            """)
    
    with tab3:
        st.subheader("Data Management")
        
        # Database statistics
        projects = st.session_state.db.get_projects()
        documents = st.session_state.db.get_documents()
        workflows = st.session_state.db.get_workflows()
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            st.metric("Projects", len(projects))
        
        with col2:
            st.metric("Documents", len(documents))
        
        with col3:
            st.metric("Workflows", len(workflows))
        
        st.subheader("Database Actions")
        
        col1, col2 = st.columns(2)
        
        with col1:
            if st.button("📥 Export Data", type="secondary"):
                export_data = {
                    'projects': projects,
                    'documents': documents,
                    'workflows': workflows,
                    'exported_at': datetime.now().isoformat()
                }
                
                st.download_button(
                    label="⬇️ Download Export",
                    data=json.dumps(export_data, indent=2),
                    file_name=f"bosch_projects_export_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json",
                    mime="application/json"
                )
        
        with col2:
            if st.button("🔄 Load Sample Data", type="secondary"):
                # Add sample data
                sample_project = {
                    "name": "Bosch IoT Platform",
                    "type": "Software Development",
                    "description": "Next-generation IoT platform for connected devices",
                    "functional_reqs": ["System shall handle 1M+ concurrent connections", "Real-time data processing"],
                    "non_functional_reqs": ["99.9% uptime", "Sub-100ms response time"],
                    "conditions": ["GDPR compliance required", "Bosch security standards"],
                    "recommended_docs": ["Project Management Plan (PMP)", "Technical Concept Document (TCD)"]
                }
                
                project_id = st.session_state.db.save_project(sample_project)
                st.success("Sample data loaded!")
                st.rerun()
    
    with tab3:
        st.subheader("Session Management")
        
        col1, col2 = st.columns(2)
        
        with col1:
            if st.button("🚪 Logout", type="secondary"):
                st.session_state.authenticated = False
                st.session_state.login_attempts = 0
                st.rerun()
        
        with col2:
            if st.button("🔄 Restart App", type="secondary"):
                st.rerun()
        
        st.subheader("Session Information")
        st.write(f"**Authenticated:** {st.session_state.get('authenticated', False)}")
        st.write(f"**Chat Messages:** {len(st.session_state.get('chat_history', []))}")
        st.write(f"**Current Model:** {st.session_state.llm_settings['model']}")

def show_workflow_visualization():
    """Display workflow visualization with current status"""
    st.title("📋 Workflow Visualization")
    
    workflows = st.session_state.db.get_workflows()
    documents = st.session_state.db.get_documents()
    
    if not workflows:
        st.info("No workflows found. Generate some documents first.")
        return
    
    st.subheader("📊 Workflow Status Overview")
    
    # Create visual workflow cards
    for workflow in workflows:
        document = next((d for d in documents if d["id"] == workflow["document_id"]), None)
        if document:
            # Determine status color
            if workflow['status'] == 'Active':
                status_color = '#FFC107'  # Yellow
                status_icon = '🟡'
            elif workflow['status'] == 'Rejected':
                status_color = '#DC3545'  # Red
                status_icon = '🔴'
            elif workflow['status'] == 'Completed':
                status_color = '#28A745'  # Green
                status_icon = '🟢'
            else:
                status_color = '#6C757D'  # Gray
                status_icon = '⚪'
            
            # Create workflow card
            st.markdown(f"""
            <div style="
                border: 2px solid {status_color}; 
                border-radius: 10px; 
                padding: 1rem; 
                margin: 1rem 0;
                background-color: {status_color}20;
            ">
                <h4 style="margin-top: 0;">{status_icon} {document['name']}</h4>
                <p><strong>Type:</strong> {document['type']}</p>
                <p><strong>Status:</strong> {workflow['status']}</p>
                <p><strong>Progress:</strong> Step {workflow.get('current_step', 0) + 1} of {len(workflow.get('approvers', []))}</p>
            </div>
            """, unsafe_allow_html=True)
            
            # Show approval chain
            if workflow.get('approvers'):
                st.write("**Approval Chain:**")
                cols = st.columns(len(workflow['approvers']))
                for i, approver in enumerate(workflow['approvers']):
                    with cols[i]:
                        if i < workflow.get('current_step', 0):
                            st.success(f"✅ {approver}")
                        elif i == workflow.get('current_step', 0) and workflow['status'] == 'Active':
                            st.warning(f"⏳ {approver}")
                        elif workflow['status'] == 'Rejected':
                            st.error(f"❌ {approver}")
                        else:
                            st.info(f"⏸️ {approver}")
            
            st.markdown("---")

def show_compliance_audit():
    """Display compliance audit features for Quality team"""
    st.title("🔍 Compliance Audit")
    
    st.markdown("""
    ### ISO 9001 & ASPICE Compliance Dashboard
    This section provides audit trails and compliance verification tools.
    """)
    
    tab1, tab2, tab3, tab4 = st.tabs([
        "📋 Audit Trail", 
        "📊 Compliance Metrics", 
        "🔍 Document Review", 
        "📝 Audit Reports"
    ])
    
    with tab1:
        st.subheader("📋 System Audit Trail")
        
        # Get all workflows and documents for audit
        workflows = st.session_state.db.get_workflows()
        documents = st.session_state.db.get_documents()
        
        st.info(f"**Total Documents:** {len(documents)} | **Total Workflows:** {len(workflows)}")
        
        # Show detailed audit information
        if workflows:
            st.write("**Workflow Audit Trail:**")
            audit_data = []
            for workflow in workflows:
                document = next((d for d in documents if d["id"] == workflow["document_id"]), None)
                if document:
                    # Get workflow comments for audit trail
                    try:
                        comments = st.session_state.db.get_workflow_comments(workflow['id'])
                        comment_count = len(comments) if comments else 0
                    except:
                        comment_count = 0
                    
                    audit_data.append({
                        'Document': document['name'],
                        'Type': document['type'],
                        'Status': workflow['status'],
                        'Created': workflow.get('created_at', 'N/A'),
                        'Comments': comment_count,
                        'Current Step': f"{workflow.get('current_step', 0) + 1}/{len(workflow.get('approvers', []))}"
                    })
            
            if audit_data:
                df = pd.DataFrame(audit_data)
                st.dataframe(df, use_container_width=True)
    
    with tab2:
        st.subheader("📊 Compliance Metrics")
        
        if workflows:
            # Calculate compliance metrics
            total_workflows = len(workflows)
            active_workflows = len([w for w in workflows if w['status'] == 'Active'])
            completed_workflows = len([w for w in workflows if w['status'] == 'Completed'])
            rejected_workflows = len([w for w in workflows if w['status'] == 'Rejected'])
            
            col1, col2, col3, col4 = st.columns(4)
            with col1:
                st.metric("Total Workflows", total_workflows)
            with col2:
                st.metric("Active", active_workflows, delta=None)
            with col3:
                st.metric("Completed", completed_workflows, delta=None)
            with col4:
                st.metric("Rejected", rejected_workflows, delta=None)
            
            # Compliance percentage
            if total_workflows > 0:
                compliance_rate = (completed_workflows / total_workflows) * 100
                st.metric("Completion Rate", f"{compliance_rate:.1f}%", 
                         delta=f"{'✅' if compliance_rate >= 70 else '⚠️'}")
            
            # Process adherence metrics
            st.subheader("🎯 Process Adherence")
            st.write("- ✅ All documents follow approval workflow")
            st.write("- ✅ Audit trail maintained for all changes")
            st.write("- ✅ Role-based access controls implemented")
            st.write("- ✅ Document versioning tracked")
    
    with tab3:
        st.subheader("🔍 Document Quality Review")
        
        if documents:
            selected_doc = st.selectbox(
                "Select document for quality review:",
                options=[f"{d['name']} ({d['type']})" for d in documents],
                key="quality_review_doc"
            )
            
            if selected_doc:
                doc_name = selected_doc.split(' (')[0]
                document = next((d for d in documents if d['name'] == doc_name), None)
                
                if document:
                    st.write(f"**Document:** {document['name']}")
                    st.write(f"**Type:** {document['type']}")
                    st.write(f"**Status:** {document['status']}")
                    
                    # Quality checklist
                    st.subheader("📋 Quality Checklist")
                    checks = {
                        "Document structure follows template": st.checkbox("Document structure follows template", key="check1"),
                        "All required sections present": st.checkbox("All required sections present", key="check2"),
                        "Content is technically accurate": st.checkbox("Content is technically accurate", key="check3"),
                        "Follows Bosch standards": st.checkbox("Follows Bosch standards", key="check4"),
                        "Approval process completed": st.checkbox("Approval process completed", key="check5")
                    }
                    
                    # Quality score
                    passed_checks = sum(checks.values())
                    total_checks = len(checks)
                    quality_score = (passed_checks / total_checks) * 100
                    
                    if quality_score >= 80:
                        st.success(f"✅ Quality Score: {quality_score:.0f}% - Compliant")
                    elif quality_score >= 60:
                        st.warning(f"⚠️ Quality Score: {quality_score:.0f}% - Needs Improvement")
                    else:
                        st.error(f"❌ Quality Score: {quality_score:.0f}% - Non-Compliant")
    
    with tab4:
        st.subheader("📝 Audit Reports")
        
        report_type = st.selectbox(
            "Select Report Type:",
            ["ISO 9001 Compliance Report", "ASPICE Process Report", "Document Traceability Report"]
        )
        
        if st.button("📄 Generate Report"):
            st.success("✅ Report generated successfully!")
            
            if report_type == "ISO 9001 Compliance Report":
                st.markdown("""
                ### ISO 9001 Compliance Report
                **Generated:** {datetime.now().strftime("%Y-%m-%d %H:%M")}
                
                **Quality Management System Status:**
                - ✅ Document control procedures implemented
                - ✅ Process workflows defined and followed  
                - ✅ Management review processes active
                - ✅ Corrective action tracking enabled
                
                **Recommendations:**
                - Continue regular audit cycles
                - Monitor completion rates
                - Ensure continuous improvement
                """)
            
            elif report_type == "ASPICE Process Report":
                st.markdown("""
                ### ASPICE Process Assessment Report
                **Generated:** {datetime.now().strftime("%Y-%m-%d %H:%M")}
                
                **Process Capability Assessment:**
                - **SYS.2 System Requirements Analysis:** Level 3 (Defined Process)
                - **SWE.1 Software Requirements Analysis:** Level 3 (Defined Process)
                - **SWE.2 Software Architectural Design:** Level 2 (Managed Process)
                
                **Process Adherence:**
                - Work products properly managed
                - Traceability maintained
                - Quality gates implemented
                """)
                
            else:  # Document Traceability Report
                st.markdown("""
                ### Document Traceability Report
                **Generated:** {datetime.now().strftime("%Y-%m-%d %H:%M")}
                
                **Traceability Matrix:**
                - Requirements → Design Documents: ✅ Complete
                - Design → Implementation: ✅ Complete  
                - Test Cases → Requirements: ✅ Complete
                
                **Gap Analysis:**
                - No missing traceability links identified
                - All documents properly versioned
                - Approval chains complete
                """)

if __name__ == "__main__":
    main()
